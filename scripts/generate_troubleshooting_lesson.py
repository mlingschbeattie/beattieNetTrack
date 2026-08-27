import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TARGET_DIR = r"c:\CustomApps\26_27_ LessonsAndAgendas"
os.makedirs(TARGET_DIR, exist_ok=True)

print("Starting generation of CompTIA 6-Step Troubleshooting Lesson & Activity...")

# ==========================================
# 1. GENERATE PRESENTATION.HTML
# ==========================================
presentation_html_path = os.path.join(TARGET_DIR, "presentation.html")

presentation_html = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Divide &amp; Conquer // CompTIA 6-Step Troubleshooting Mastery</title>
  
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=IBM+Plex+Sans+Arabic:wght@400;600;700&family=Inter:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;500;700;800&family=Space+Grotesk:wght@600;700&display=swap" rel="stylesheet">
  
  <style>
    :root {
      --bg: #040608;
      --panel: #090e13;
      --panel-elevated: #101720;
      --line: #1b272c;
      --line-bright: #2d424b;
      --green: #39ff9e;
      --green-glow: rgba(57, 255, 158, 0.35);
      --green-dim: #134e35;
      --cyan: #38bdf8;
      --cyan-glow: rgba(56, 189, 248, 0.35);
      --amber: #ffb454;
      --red: #ff5c72;
      --purple: #a855f7;
      --text: #e2f4ea;
      --text-muted: #829a90;
      --text-dim: #4d665c;
    }

    * { box-sizing: border-box; margin: 0; padding: 0; }

    body {
      background-color: var(--bg);
      color: var(--text);
      font-family: 'JetBrains Mono', monospace;
      min-height: 100vh;
      overflow: hidden;
      display: flex;
      flex-direction: column;
      user-select: none;
      position: relative;
    }

    body::before {
      content: "";
      position: fixed;
      inset: 0;
      background-image:
        linear-gradient(rgba(57, 255, 158, 0.04) 1px, transparent 1px),
        linear-gradient(90deg, rgba(57, 255, 158, 0.04) 1px, transparent 1px);
      background-size: 40px 40px;
      pointer-events: none;
      z-index: 1;
    }

    body::after {
      content: "";
      position: fixed;
      inset: 0;
      background: 
        repeating-linear-gradient(to bottom, rgba(0,0,0,0) 0px, rgba(0,0,0,0) 2px, rgba(0,0,0,0.18) 3px, rgba(0,0,0,0) 4px),
        radial-gradient(ellipse at center, transparent 50%, rgba(0,0,0,0.7) 100%);
      pointer-events: none;
      mix-blend-mode: multiply;
      z-index: 2;
    }

    :root:lang(ar) {
      font-family: 'IBM Plex Sans Arabic', 'Inter', sans-serif;
    }

    [dir="rtl"] {
      text-align: start;
    }

    [dir="rtl"] code,
    [dir="rtl"] pre,
    [dir="rtl"] .font-mono,
    [dir="rtl"] .cli-box {
      direction: ltr !important;
      unicode-bidi: isolate;
      text-align: left;
    }

    .topbar {
      position: relative;
      z-index: 10;
      display: flex;
      justify-content: space-between;
      align-items: center;
      padding: 0.75rem 2rem;
      border-bottom: 1px solid var(--line);
      background: rgba(9, 14, 19, 0.85);
      backdrop-filter: blur(8px);
      font-size: 0.8rem;
      color: var(--text-muted);
    }

    .brand {
      display: flex;
      align-items: center;
      gap: 0.6rem;
      font-weight: 700;
      letter-spacing: 0.08em;
      color: #fff;
    }
    .brand-dot { color: var(--green); animation: pulse 2s infinite; }

    .nav-controls {
      display: flex;
      align-items: center;
      gap: 0.8rem;
    }

    .btn-ctrl {
      background: var(--panel-elevated);
      border: 1px solid var(--line-bright);
      color: var(--text);
      padding: 0.35rem 0.8rem;
      border-radius: 4px;
      font-family: inherit;
      font-size: 0.75rem;
      font-weight: 600;
      cursor: pointer;
      display: inline-flex;
      align-items: center;
      gap: 0.4rem;
      transition: all 0.15s ease;
      text-decoration: none;
    }
    .btn-ctrl:hover {
      border-color: var(--green);
      color: var(--green);
      box-shadow: 0 0 10px var(--green-glow);
    }
    .btn-ctrl.active {
      background: var(--green-dim);
      border-color: var(--green);
      color: var(--green);
    }

    .slide-counter {
      font-weight: 800;
      color: var(--green);
      font-size: 0.85rem;
      letter-spacing: 0.05em;
      padding: 0 0.5rem;
    }

    .stage {
      position: relative;
      z-index: 5;
      flex: 1;
      display: flex;
      align-items: center;
      justify-content: center;
      padding: 1.5rem 3rem;
      overflow: hidden;
    }

    .slide {
      position: absolute;
      width: 100%;
      max-width: 1350px;
      height: 82vh;
      max-height: 850px;
      background: var(--panel);
      border: 1px solid var(--line-bright);
      border-radius: 12px;
      padding: 2.2rem 3rem;
      display: flex;
      flex-direction: column;
      box-shadow: 0 20px 50px rgba(0,0,0,0.8), 0 0 30px rgba(57, 255, 158, 0.03);
      opacity: 0;
      transform: translateX(60px) scale(0.97);
      pointer-events: none;
      transition: opacity 0.35s cubic-bezier(0.16, 1, 0.3, 1), transform 0.35s cubic-bezier(0.16, 1, 0.3, 1);
      overflow-y: auto;
    }

    .slide.active {
      opacity: 1;
      transform: translateX(0) scale(1);
      pointer-events: auto;
    }

    .slide.prev {
      opacity: 0;
      transform: translateX(-60px) scale(0.97);
    }

    .slide-badge {
      display: inline-flex;
      align-items: center;
      gap: 0.5rem;
      font-size: 0.75rem;
      font-weight: 800;
      letter-spacing: 0.15em;
      text-transform: uppercase;
      color: var(--green);
      background: rgba(57, 255, 158, 0.08);
      border: 1px solid var(--green-dim);
      padding: 0.3rem 0.8rem;
      border-radius: 4px;
      width: fit-content;
      margin-bottom: 0.8rem;
    }

    .slide-badge.cyan { color: var(--cyan); background: rgba(56, 189, 248, 0.08); border-color: rgba(56, 189, 248, 0.3); }
    .slide-badge.amber { color: var(--amber); background: rgba(255, 180, 84, 0.08); border-color: rgba(255, 180, 84, 0.3); }
    .slide-badge.purple { color: var(--purple); background: rgba(168, 85, 247, 0.08); border-color: rgba(168, 85, 247, 0.3); }
    .slide-badge.red { color: var(--red); background: rgba(255, 92, 114, 0.08); border-color: rgba(255, 92, 114, 0.3); }

    .slide-title {
      font-family: 'Space Grotesk', sans-serif;
      font-size: clamp(1.8rem, 3vw, 2.6rem);
      font-weight: 700;
      color: #fff;
      line-height: 1.15;
      margin-bottom: 0.4rem;
      letter-spacing: -0.02em;
    }

    .slide-subtitle {
      font-size: 1rem;
      color: var(--text-muted);
      margin-bottom: 1.8rem;
      line-height: 1.5;
    }

    .grid-2 {
      display: grid;
      grid-template-columns: 1fr 1fr;
      gap: 1.6rem;
      flex: 1;
    }

    .grid-3 {
      display: grid;
      grid-template-columns: 1fr 1fr 1fr;
      gap: 1.4rem;
      flex: 1;
    }

    .grid-6 {
      display: grid;
      grid-template-columns: repeat(3, 1fr);
      grid-template-rows: repeat(2, 1fr);
      gap: 1.2rem;
      flex: 1;
    }

    .card {
      background: var(--panel-elevated);
      border: 1px solid var(--line-bright);
      border-radius: 8px;
      padding: 1.3rem 1.5rem;
      position: relative;
      display: flex;
      flex-direction: column;
      transition: all 0.2s ease;
    }
    .card:hover {
      border-color: var(--green);
      transform: translateY(-2px);
      box-shadow: 0 8px 24px rgba(0,0,0,0.5);
    }

    .card-num {
      font-size: 1.6rem;
      font-weight: 800;
      color: var(--green);
      margin-bottom: 0.4rem;
      font-family: 'Space Grotesk', sans-serif;
    }
    .card-num.cyan { color: var(--cyan); }
    .card-num.amber { color: var(--amber); }
    .card-num.red { color: var(--red); }
    .card-num.purple { color: var(--purple); }

    .card-title {
      font-size: 1.1rem;
      font-weight: 700;
      color: #fff;
      margin-bottom: 0.5rem;
      font-family: 'Space Grotesk', sans-serif;
    }

    .card-desc {
      font-size: 0.85rem;
      color: var(--text-muted);
      line-height: 1.6;
    }

    .card-desc strong { color: #fff; }

    .card-tag {
      margin-top: auto;
      padding-top: 0.8rem;
      font-size: 0.72rem;
      color: var(--green);
      letter-spacing: 0.08em;
      text-transform: uppercase;
      font-weight: 700;
    }

    .cli-box {
      background: #020305;
      border: 1px solid var(--line);
      border-left: 3px solid var(--green);
      padding: 0.8rem 1.1rem;
      border-radius: 4px;
      font-family: 'JetBrains Mono', monospace;
      font-size: 0.82rem;
      color: var(--green);
      margin: 0.7rem 0;
      white-space: pre-wrap;
    }

    .cli-box.amber { border-left-color: var(--amber); color: var(--amber); }
    .cli-box.cyan { border-left-color: var(--cyan); color: var(--cyan); }
    .cli-box.red { border-left-color: var(--red); color: var(--red); }

    .sim-container {
      background: #020305;
      border: 1px solid var(--line-bright);
      border-radius: 8px;
      padding: 1.4rem;
      display: flex;
      flex-direction: column;
      gap: 1rem;
    }

    .osi-ladder {
      display: flex;
      flex-direction: column;
      gap: 0.5rem;
    }

    .osi-step {
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 0.65rem 1.1rem;
      background: var(--panel-elevated);
      border: 1px solid var(--line);
      border-radius: 6px;
      cursor: pointer;
      transition: all 0.2s ease;
    }
    .osi-step:hover {
      border-color: var(--cyan);
      background: #141e2b;
    }
    .osi-step.selected {
      border-color: var(--green);
      background: rgba(57, 255, 158, 0.1);
      box-shadow: 0 0 15px rgba(57, 255, 158, 0.2);
    }
    .osi-step.failed {
      border-color: var(--red);
      background: rgba(255, 92, 114, 0.1);
    }

    .osi-step-left {
      display: flex;
      align-items: center;
      gap: 0.9rem;
    }
    .osi-num {
      font-weight: 800;
      color: var(--cyan);
      font-size: 0.85rem;
      width: 2rem;
    }
    .osi-name {
      font-weight: 700;
      color: #fff;
      font-size: 0.9rem;
    }
    .osi-cmd {
      font-size: 0.75rem;
      color: var(--amber);
    }

    .osi-status {
      font-weight: 800;
      font-size: 0.72rem;
      padding: 0.2rem 0.55rem;
      border-radius: 3px;
      text-transform: uppercase;
    }
    .osi-status.pass { background: rgba(57, 255, 158, 0.2); color: var(--green); }
    .osi-status.fail { background: rgba(255, 92, 114, 0.2); color: var(--red); }
    .osi-status.untested { background: rgba(255, 255, 255, 0.05); color: var(--text-muted); }

    .bottombar {
      position: relative;
      z-index: 10;
      display: flex;
      justify-content: space-between;
      align-items: center;
      padding: 0.8rem 2rem;
      border-top: 1px solid var(--line);
      background: rgba(9, 14, 19, 0.95);
      backdrop-filter: blur(8px);
    }

    .progress-track {
      flex: 1;
      height: 4px;
      background: var(--line);
      margin: 0 2rem;
      border-radius: 2px;
      overflow: hidden;
    }
    .progress-fill {
      height: 100%;
      background: linear-gradient(90deg, var(--cyan), var(--green));
      width: 6.66%;
      transition: width 0.3s ease;
      box-shadow: 0 0 8px var(--green-glow);
    }

    .key-hints {
      font-size: 0.75rem;
      color: var(--text-dim);
      display: flex;
      gap: 1rem;
    }
    .key-hints span kbd {
      background: var(--panel-elevated);
      border: 1px solid var(--line-bright);
      padding: 0.15rem 0.4rem;
      border-radius: 3px;
      color: var(--text-muted);
      font-weight: 700;
    }

    .presenter-drawer {
      position: fixed;
      bottom: 60px;
      left: 2rem;
      right: 2rem;
      max-height: 220px;
      background: #020305;
      border: 1px solid var(--amber);
      border-radius: 8px;
      padding: 1.2rem 1.6rem;
      z-index: 100;
      box-shadow: 0 10px 40px rgba(0,0,0,0.9);
      display: none;
      overflow-y: auto;
    }
    .presenter-drawer.open { display: block; }
    .presenter-header {
      display: flex;
      justify-content: space-between;
      align-items: center;
      margin-bottom: 0.6rem;
      color: var(--amber);
      font-weight: 800;
      font-size: 0.8rem;
      letter-spacing: 0.1em;
      text-transform: uppercase;
    }
    .presenter-text {
      font-size: 0.9rem;
      line-height: 1.6;
      color: #fed7aa;
    }

    @keyframes pulse {
      0%, 100% { opacity: 1; transform: scale(1); }
      50% { opacity: 0.4; transform: scale(0.85); }
    }
  </style>
</head>
<body>

  <header class="topbar">
    <div class="brand">
      <span class="brand-dot">●</span>
      <span>BEATTIE-NET // CYBER &amp; NET ENGINEERING</span>
      <span style="color: var(--text-dim);">|</span>
      <span style="color: var(--cyan); font-weight: 600;">COMPTIA 6-STEP TROUBLESHOOTING</span>
    </div>

    <div class="nav-controls">
      <div style="display:flex; background:var(--panel-elevated); border:1px solid var(--line-bright); border-radius:4px; padding:2px; font-size:0.75rem;">
        <button id="lang-en" class="btn-ctrl active" onclick="setLanguage('en')" style="border:none; padding:0.2rem 0.6rem;">EN</button>
        <button id="lang-ar" class="btn-ctrl" onclick="setLanguage('ar')" style="border:none; padding:0.2rem 0.6rem;">عربي</button>
        <button id="lang-uk" class="btn-ctrl" onclick="setLanguage('uk')" style="border:none; padding:0.2rem 0.6rem;">УКР</button>
      </div>

      <button id="btn-notes" class="btn-ctrl" onclick="toggleNotes()" title="Toggle Instructor Teaching Notes (Press 'N')">
        📝 <span class="txt-notes">Notes [N]</span>
      </button>
      
      <button id="btn-fs" class="btn-ctrl" onclick="toggleFullScreen()" title="Fullscreen Mode (Press 'F')">
        ⛶ <span class="txt-fs">Fullscreen [F]</span>
      </button>

      <a href="activity.html" class="btn-ctrl" style="border-color:var(--green); color:var(--green);">
        ⚡ <span class="txt-act">Launch Lab Activity →</span>
      </a>
    </div>
  </header>

  <main class="stage" id="stage">

    <!-- SLIDE 1 -->
    <section class="slide active" data-slide="1">
      <div class="slide-badge">LESSON // COMPTIA CORE METHODOLOGY</div>
      <h1 class="slide-title">Divide &amp; Conquer: The 6-Step Model</h1>
      <p class="slide-subtitle">Systematic Fault Isolation, Binary Search Troubleshooting, &amp; Enterprise Triage</p>

      <div class="grid-2">
        <div class="card" style="border-left: 4px solid var(--green);">
          <div class="card-num">01</div>
          <h2 class="card-title">Stop Guessing. Start Isolating.</h2>
          <p class="card-desc">
            Amateur technicians guess and swap random cables. Enterprise network engineers and ethical hackers use <strong>formal deductive models</strong>. Today, you master the exact 6-step framework codified by <strong>CompTIA A+, Network+, and Security+</strong>.
          </p>
          <div class="cli-box">root@beattie-net:~# ./isolate_fault.sh --target="enterprise-network"
[+] Establishing baseline telemetry...
[+] Initializing 6-step diagnostic protocol...</div>
          <div class="card-tag">CompTIA A+ (220-1101/1102) &middot; Network+ (N10-008/009)</div>
        </div>

        <div class="card" style="border-left: 4px solid var(--cyan);">
          <div class="card-num cyan">02</div>
          <h2 class="card-title">The Binary Search Algorithm</h2>
          <p class="card-desc">
            Linear troubleshooting takes $O(N)$ time. By applying <strong>Divide &amp; Conquer at the OSI midpoint</strong> (Layer 3/4), you eliminate 50% of the entire problem space in a single CLI command (`ping 127.0.0.1` &rarr; `ping gateway` &rarr; `ping 8.8.8.8` &rarr; `nslookup`).
          </p>
          <div class="cli-box cyan">Algorithm Efficiency:
7 Linear Tests -> 100% Time
Divide & Conquer (3 Tests) -> 57% Faster Mean Time to Resolution (MTTR)</div>
          <div class="card-tag">Mean Time to Resolution (MTTR) Optimization</div>
        </div>
      </div>
    </section>

    <!-- SLIDE 2 -->
    <section class="slide" data-slide="2">
      <div class="slide-badge amber">INDUSTRY REALITY CHECK</div>
      <h1 class="slide-title">The High Cost of Random Guessing</h1>
      <p class="slide-subtitle">Why systematic diagnosis is mandatory in corporate and cybersecurity environments</p>

      <div class="grid-3">
        <div class="card" style="border-top: 3px solid var(--red);">
          <div class="card-num red">01</div>
          <h2 class="card-title">Wasted Downtime ($$$)</h2>
          <p class="card-desc">
            In enterprise datacenters, network downtime averages <strong>$5,600 per minute</strong>. Swapping parts randomly without a verified hypothesis creates multi-hour outages that cost companies millions.
          </p>
          <div class="card-tag" style="color:var(--red);">Financial &amp; Operational Risk</div>
        </div>

        <div class="card" style="border-top: 3px solid var(--amber);">
          <div class="card-num amber">02</div>
          <h2 class="card-title">Collateral Damage</h2>
          <p class="card-desc">
            Modifying settings, rebooting production routers, or updating drivers without an action plan often breaks <strong>unrelated dependent services</strong>, multiplying one ticket into fifty.
          </p>
          <div class="card-tag" style="color:var(--amber);">Configuration Drift &amp; Chaos</div>
        </div>

        <div class="card" style="border-top: 3px solid var(--green);">
          <div class="card-num">03</div>
          <h2 class="card-title">Un-reproducible Fixes</h2>
          <p class="card-desc">
            If you tweak six settings simultaneously and the network starts working, <strong>which one fixed it?</strong> You don't know, and when it breaks next week, you will start from square one.
          </p>
          <div class="card-tag">Scientific Methodology</div>
        </div>
      </div>
    </section>

    <!-- SLIDE 3 -->
    <section class="slide" data-slide="3">
      <div class="slide-badge cyan">THE COMPTIA 6-STEP STANDARD</div>
      <h1 class="slide-title">The Complete Troubleshooting Lifecycle</h1>
      <p class="slide-subtitle">The industry-standard sequential methodology required for CompTIA Certification</p>

      <div class="grid-6">
        <div class="card">
          <div class="card-num">1</div>
          <h3 class="card-title">Identify the Problem</h3>
          <p class="card-desc">Question the user, identify changes, review system logs, duplicate the symptom.</p>
        </div>
        <div class="card">
          <div class="card-num cyan">2</div>
          <h3 class="card-title">Establish a Theory</h3>
          <p class="card-desc">Question the obvious. Consider Top-to-Bottom, Bottom-to-Top, or Divide &amp; Conquer.</p>
        </div>
        <div class="card">
          <div class="card-num amber">3</div>
          <h3 class="card-title">Test the Theory</h3>
          <p class="card-desc">Confirm root cause. If disproven, form a new theory or escalate immediately.</p>
        </div>
        <div class="card">
          <div class="card-num purple">4</div>
          <h3 class="card-title">Plan &amp; Implement</h3>
          <p class="card-desc">Establish action plan, identify potential side effects, and safely execute solution.</p>
        </div>
        <div class="card">
          <div class="card-num">5</div>
          <h3 class="card-title">Verify &amp; Prevent</h3>
          <p class="card-desc">Verify full functionality with user; implement proactive preventive measures.</p>
        </div>
        <div class="card">
          <div class="card-num cyan">6</div>
          <h3 class="card-title">Document Everything</h3>
          <p class="card-desc">Record findings, actions, and outcomes in enterprise ticketing &amp; knowledge base.</p>
        </div>
      </div>
    </section>

    <!-- SLIDE 4 -->
    <section class="slide" data-slide="4">
      <div class="slide-badge">STEP 1 IN-DEPTH</div>
      <h1 class="slide-title">Step 1: Identify the Problem</h1>
      <p class="slide-subtitle">Gathering raw telemetry before touching a single configuration file</p>

      <div class="grid-2">
        <div class="card">
          <h2 class="card-title">Interrogation &amp; Change Tracking</h2>
          <p class="card-desc">
            Users rarely report the technical root cause; they report symptoms (e.g. <em>"The server is dead"</em>). Your job as a Tier 2 technician:
          </p>
          <ul style="margin: 0.8rem 0 0 1.2rem; color: var(--text-muted); font-size: 0.88rem; line-height: 1.8;">
            <li><strong>User Questioning:</strong> When did it start? What exactly were you doing?</li>
            <li><strong>Environmental Changes:</strong> Was new software installed? Was an OS update pushed overnight? Did physical cables move?</li>
            <li><strong>Duplication:</strong> Can you replicate the exact failure on command?</li>
            <li><strong>Scope Determination:</strong> Is this 1 workstation, 1 VLAN, or the entire building?</li>
          </ul>
        </div>

        <div class="card">
          <h2 class="card-title">Telemetry &amp; Log Extraction</h2>
          <p class="card-desc">
            Always inspect authoritative system telemetry before theorizing:
          </p>
          <div class="cli-box"># Windows PowerShell Diagnostic Commands:
Get-EventLog -LogName System -Newest 10 -EntryType Error
ipconfig /all
Test-NetConnection -ComputerName 10.0.0.1 -Port 443

# Linux Forensic Inspection:
journalctl -xe -u networking.service
dmesg | grep -E "eth|wlan|link"</div>
        </div>
      </div>
    </section>

    <!-- SLIDE 5 -->
    <section class="slide" data-slide="5">
      <div class="slide-badge cyan">STEPS 2 &amp; 3 IN-DEPTH</div>
      <h1 class="slide-title">Step 2 &amp; 3: Theorize &amp; Falsify</h1>
      <p class="slide-subtitle">Hypothesis generation and rapid scientific testing</p>

      <div class="grid-2">
        <div class="card">
          <div class="card-num cyan">02</div>
          <h2 class="card-title">Establish Theory of Probable Cause</h2>
          <p class="card-desc">
            <strong>Question the Obvious:</strong> Is the patch cable disconnected? Is the switch port unassigned? Is caps lock on? Did the DHCP lease expire?
          </p>
          <p class="card-desc" style="margin-top: 0.8rem;">
            Rank your theories by likelihood. Always test the least invasive, highest-probability hypothesis first.
          </p>
          <div class="cli-box cyan">Theory A (Highest Prob): Local DNS resolver timeout
Theory B (Medium Prob): Default gateway dropped ARP
Theory C (Low Prob): ISP fiber cut across town</div>
        </div>

        <div class="card">
          <div class="card-num amber">03</div>
          <h2 class="card-title">Test the Theory (Falsification)</h2>
          <p class="card-desc">
            Execute targeted, isolated tests. A good test produces a clean binary outcome: <strong>CONFIRMED</strong> or <strong>DISPROVEN</strong>.
          </p>
          <p class="card-desc" style="margin-top: 0.8rem;">
            <strong>Critical CompTIA Rule:</strong> If your test disproves the theory, <em>do not execute a repair plan</em>. Form a new theory or escalate to Tier 3 / Network Architecture.
          </p>
          <div class="cli-box amber">C:\> ping 8.8.8.8  -> Reply from 8.8.8.8 (PASS!)
C:\> nslookup google.com -> DNS request timed out (FAIL!)
[+] THEORY CONFIRMED: IP routing works; DNS resolution failed.</div>
        </div>
      </div>
    </section>

    <!-- SLIDE 6 -->
    <section class="slide" data-slide="6">
      <div class="slide-badge purple">STEP 4 IN-DEPTH</div>
      <h1 class="slide-title">Step 4: Plan of Action &amp; Implementation</h1>
      <p class="slide-subtitle">Safe execution with change management, risk mitigation, &amp; rollback plans</p>

      <div class="grid-2">
        <div class="card">
          <h2 class="card-title">Constructing the Action Plan</h2>
          <p class="card-desc">
            Never deploy a fix blind. A professional action plan consists of three core pillars:
          </p>
          <ul style="margin: 0.8rem 0 0 1.2rem; color: var(--text-muted); font-size: 0.88rem; line-height: 1.8;">
            <li><strong>Step-by-Step Remediation:</strong> Exact commands, config syntax, and files to modify.</li>
            <li><strong>Potential Side Effects:</strong> Will bouncing this switch interface drop active VoIP phone calls in room 204?</li>
            <li><strong>Rollback Protocol:</strong> If the update corrupts the database, how do you restore the exact previous snapshot in under 5 minutes?</li>
          </ul>
        </div>

        <div class="card">
          <h2 class="card-title">Change Windows &amp; Authorization</h2>
          <div class="cli-box purple">ENTERPRISE CHANGE TICKET #8849
Scope: Core Switch VLAN 20 IP Helper Update
Outage Window: 05:00 - 05:15 EST (Low impact)
Impact Analysis: 45 Workstations will renew DHCP leases
Rollback Plan: Restore running-config from TFTP /backup/sw01_prev.cfg</div>
          <p class="card-desc" style="margin-top: 0.8rem;">
            Always obtain appropriate authorization (Change Advisory Board / Lead Systems Admin) before modifying enterprise infrastructure.
          </p>
        </div>
      </div>
    </section>

    <!-- SLIDE 7 -->
    <section class="slide" data-slide="7">
      <div class="slide-badge">STEPS 5 &amp; 6 IN-DEPTH</div>
      <h1 class="slide-title">Step 5 &amp; 6: Verification &amp; Documentation</h1>
      <p class="slide-subtitle">Closing the loop, preventative hardening, and building organizational memory</p>

      <div class="grid-2">
        <div class="card">
          <div class="card-num">05</div>
          <h2 class="card-title">Verify &amp; Prevent</h2>
          <p class="card-desc">
            <strong>1. Full System Verification:</strong> Do not just check your command line; have the end user perform their real daily workflow in front of you.
          </p>
          <p class="card-desc" style="margin-top: 0.8rem;">
            <strong>2. Preventative Hardening:</strong> Why did it fail? If a cable failed from being stepped on, install floor raceways. If an IP address conflicted, configure DHCP reservation and Dynamic ARP Inspection (DAI).
          </p>
        </div>

        <div class="card">
          <div class="card-num cyan">06</div>
          <h2 class="card-title">Document Findings, Actions &amp; Outcomes</h2>
          <p class="card-desc">
            If you don't document it, the problem was never solved. Standard ticketing format:
          </p>
          <div class="cli-box">TICKET RESOLUTION REPORT
Symptom: Lab Workstation 12 showing 169.254.x.x (APIPA)
Root Cause: DHCP Scope 192.168.10.0/24 100% exhausted
Remediation: Expanded subnet to /23, reduced lease time to 4h
Verification: Workstation acquired 192.168.10.142; pinged Gateway
Preventative Action: Configured 80% capacity alert on DHCP server</div>
        </div>
      </div>
    </section>

    <!-- SLIDE 8 -->
    <section class="slide" data-slide="8">
      <div class="slide-badge cyan">CORE STRATEGY</div>
      <h1 class="slide-title">The "Divide &amp; Conquer" Binary Search</h1>
      <p class="slide-subtitle">Slashing diagnosis time by testing the midpoint of the OSI 7-Layer Stack</p>

      <div class="grid-2">
        <div class="card" style="border-left: 4px solid var(--red);">
          <h2 class="card-title" style="color:var(--red);">The Slow Way: Linear OSI Walk (7 Steps)</h2>
          <p class="card-desc">
            Testing Layer 1 (Cable) &rarr; Layer 2 (Switch MAC) &rarr; Layer 3 (IP Gateway) &rarr; Layer 4 (TCP Port) &rarr; Layer 5 (Session) &rarr; Layer 6 (SSL) &rarr; Layer 7 (App).
          </p>
          <div class="cli-box red">Total Sequential Tests: 7
Time Required: 100%
Frustration Factor: HIGH</div>
        </div>

        <div class="card" style="border-left: 4px solid var(--green);">
          <h2 class="card-title" style="color:var(--green);">The Pro Way: Midpoint Binary Search (3 Steps)</h2>
          <p class="card-desc">
            Test directly at <strong>Layer 3 (Network Layer)</strong>.
          </p>
          <p class="card-desc" style="margin-top: 0.6rem;">
            &bull; <strong>If Layer 3 PASSES:</strong> Layers 1, 2, and 3 are 100% confirmed healthy! Never touch a patch cable or switch port. Zoom straight to Layer 4-7.
          </p>
          <p class="card-desc" style="margin-top: 0.6rem;">
            &bull; <strong>If Layer 3 FAILS:</strong> The issue is definitely in Layer 1-3. Never waste time checking browser settings, firewalls, or TLS certs.
          </p>
        </div>
      </div>
    </section>

    <!-- SLIDE 9 -->
    <section class="slide" data-slide="9">
      <div class="slide-badge amber">THE 5-STEP MIDPOINT LADDER</div>
      <h1 class="slide-title">The Golden Diagnostic Sequence</h1>
      <p class="slide-subtitle">Memorize these 5 sequential CLI commands for enterprise network troubleshooting</p>

      <div style="display:flex; flex-direction:column; gap:0.7rem; flex:1;">
        <div class="card" style="padding:0.7rem 1.1rem; flex-direction:row; align-items:center; justify-content:space-between;">
          <div><strong style="color:var(--cyan);">1. Test Local TCP/IP Stack:</strong> <code>ping 127.0.0.1</code> / <code>ping ::1</code></div>
          <span style="color:var(--text-muted); font-size:0.8rem;">Verifies local NIC driver &amp; OS protocol stack</span>
        </div>
        <div class="card" style="padding:0.7rem 1.1rem; flex-direction:row; align-items:center; justify-content:space-between;">
          <div><strong style="color:var(--green);">2. Test Local Gateway (LAN Exit):</strong> <code>ping 192.168.1.1</code></div>
          <span style="color:var(--text-muted); font-size:0.8rem;">Verifies Patch Cable, Wall Jack, Switch, &amp; Router Interface (L1-L3)</span>
        </div>
        <div class="card" style="padding:0.7rem 1.1rem; flex-direction:row; align-items:center; justify-content:space-between;">
          <div><strong style="color:var(--amber);">3. Test External Public IP (WAN):</strong> <code>ping 8.8.8.8</code> / <code>1.1.1.1</code></div>
          <span style="color:var(--text-muted); font-size:0.8rem;">Verifies ISP modem, NAT, Default Route, &amp; WAN Internet Routing</span>
        </div>
        <div class="card" style="padding:0.7rem 1.1rem; flex-direction:row; align-items:center; justify-content:space-between;">
          <div><strong style="color:var(--purple);">4. Test DNS Name Resolution:</strong> <code>nslookup google.com</code> / <code>ping google.com</code></div>
          <span style="color:var(--text-muted); font-size:0.8rem;">Verifies Port 53 DNS resolver, root hints, &amp; caching server</span>
        </div>
        <div class="card" style="padding:0.7rem 1.1rem; flex-direction:row; align-items:center; justify-content:space-between;">
          <div><strong style="color:var(--red);">5. Test Application Port &amp; Service:</strong> <code>Test-NetConnection -Port 443</code> / <code>curl -I</code></div>
          <span style="color:var(--text-muted); font-size:0.8rem;">Verifies Web Daemon, TLS Handshake, &amp; State Inspection Firewall</span>
        </div>
      </div>
    </section>

    <!-- SLIDE 10 -->
    <section class="slide" data-slide="10">
      <div class="slide-badge cyan">INTERACTIVE LAB DEMO</div>
      <h1 class="slide-title">Live Midpoint Diagnostic Sandbox</h1>
      <p class="slide-subtitle">Click each diagnostic test to see how the problem space splits in real-time</p>

      <div class="grid-2">
        <div class="sim-container">
          <div style="font-size:0.85rem; color:var(--text-muted); font-weight:700;">SELECT DIAGNOSTIC COMMAND TO EXECUTE:</div>
          <div class="osi-ladder">
            <div class="osi-step" onclick="runSim(1)">
              <div class="osi-step-left">
                <span class="osi-num">01</span>
                <div><div class="osi-name">Loopback Adapter</div><div class="osi-cmd">ping 127.0.0.1</div></div>
              </div>
              <span id="sim-status-1" class="osi-status untested">READY</span>
            </div>

            <div class="osi-step" onclick="runSim(2)">
              <div class="osi-step-left">
                <span class="osi-num">02</span>
                <div><div class="osi-name">Default Gateway (LAN)</div><div class="osi-cmd">ping 192.168.1.1</div></div>
              </div>
              <span id="sim-status-2" class="osi-status untested">READY</span>
            </div>

            <div class="osi-step" onclick="runSim(3)">
              <div class="osi-step-left">
                <span class="osi-num">03</span>
                <div><div class="osi-name">External Public IP (WAN)</div><div class="osi-cmd">ping 8.8.8.8</div></div>
              </div>
              <span id="sim-status-3" class="osi-status untested">READY</span>
            </div>

            <div class="osi-step" onclick="runSim(4)">
              <div class="osi-step-left">
                <span class="osi-num">04</span>
                <div><div class="osi-name">Domain Name (DNS)</div><div class="osi-cmd">nslookup google.com</div></div>
              </div>
              <span id="sim-status-4" class="osi-status untested">READY</span>
            </div>
          </div>
        </div>

        <div class="card" style="justify-content:flex-start;">
          <h2 class="card-title">Live Diagnostic Telemetry</h2>
          <div id="sim-output" class="cli-box" style="height: 220px; overflow-y: auto;">Click any command on the left to initiate the binary search protocol...</div>
          <div id="sim-verdict" style="margin-top: 0.8rem; font-size: 0.88rem; color: var(--green); font-weight: 700;"></div>
        </div>
      </div>
    </section>

    <!-- SLIDE 11 -->
    <section class="slide" data-slide="11">
      <div class="slide-badge amber">CASE STUDY 1 // TIER 1 TRIAGE</div>
      <h1 class="slide-title">The 169.254.x.x (APIPA) Outage</h1>
      <p class="slide-subtitle">Symptom: 8 lab workstations suddenly display "No Internet Access"</p>

      <div class="grid-2">
        <div class="card">
          <h2 class="card-title">The Incident Scenario</h2>
          <p class="card-desc">
            Students boot up lab computers in Room 104. Link lights on NICs are solid green. However, running <code>ipconfig</code> reveals:
          </p>
          <div class="cli-box amber">IPv4 Address. . . . . . . . . . . : 169.254.144.82
Subnet Mask . . . . . . . . . . . : 255.255.0.0
Default Gateway . . . . . . . . . : [BLANK]</div>
          <p class="card-desc">
            <strong>Automatic Private IP Addressing (APIPA / RFC 3927):</strong> Occurs when the Windows DHCP client sends a DHCPDISCOVER broadcast but receives 0 responses.
          </p>
        </div>

        <div class="card">
          <h2 class="card-title">Divide &amp; Conquer Action Flow</h2>
          <ul style="margin: 0.6rem 0 0 1.2rem; color: var(--text-muted); font-size: 0.85rem; line-height: 1.8;">
            <li><strong>Step 1:</strong> Link lights on &rarr; Physical Layer 1 &amp; Layer 2 link are UP.</li>
            <li><strong>Step 2:</strong> Theory &rarr; Windows DHCP Server service crashed, or switch VLAN IP-Helper address was removed.</li>
            <li><strong>Step 3:</strong> Test &rarr; Check DHCP server daemon status via <code>Get-Service dhcpserver</code>.</li>
            <li><strong>Step 4:</strong> Fix &rarr; Restart DHCP daemon &amp; expand scope.</li>
            <li><strong>Step 5:</strong> Verify &rarr; Run <code>ipconfig /renew</code> on client; verify address <code>192.168.10.45</code> received.</li>
            <li><strong>Step 6:</strong> Document &rarr; Log DHCP lease pool expansion in ticketing system.</li>
          </ul>
        </div>
      </div>
    </section>

    <!-- SLIDE 12 -->
    <section class="slide" data-slide="12">
      <div class="slide-badge red">CASE STUDY 2 // CYBERSECURITY TIER 2</div>
      <h1 class="slide-title">Suspicious Outbound Port 22 Exfiltration</h1>
      <p class="slide-subtitle">Symptom: Intrusion Detection System (IDS) alerts on high-bandwidth outbound SSH traffic</p>

      <div class="grid-2">
        <div class="card" style="border-left: 4px solid var(--red);">
          <h2 class="card-title">The Security Incident</h2>
          <p class="card-desc">
            Snort/Suricata IDS triggers high-severity alert: <strong>Workstation-04 (192.168.10.84)</strong> is sending 2.4 GB of encrypted traffic to a known malicious Russian VPS on Port 22 (SSH).
          </p>
          <div class="cli-box red">ALERT: [1:2010935:2] ET EXPLOIT Outbound SSH Connection to Threat Intel C2 IP 185.220.101.5
Src: 192.168.10.84:54210 -> Dst: 185.220.101.5:22</div>
          <p class="card-desc" style="margin-top:0.6rem;">
            <strong>Immediate Security Action:</strong> Isolate the compromised host from the physical network before the attacker executes lateral movement.
          </p>
        </div>

        <div class="card">
          <h2 class="card-title">Incident Response 6-Step Triage</h2>
          <ul style="margin: 0.6rem 0 0 1.2rem; color: var(--text-muted); font-size: 0.85rem; line-height: 1.8;">
            <li><strong>Step 1:</strong> Identify compromised host &amp; quarantine endpoint (disable switch port).</li>
            <li><strong>Step 2:</strong> Theorize unauthorized persistence / reverse shell script.</li>
            <li><strong>Step 3:</strong> Inspect processes with <code>netstat -ano | findstr :22</code> and memory dump.</li>
            <li><strong>Step 4:</strong> Terminate rogue PID, revoke compromised credentials, block C2 IP on firewall.</li>
            <li><strong>Step 5:</strong> Full forensic antivirus scan, re-image system from clean golden image.</li>
            <li><strong>Step 6:</strong> File Incident Response report &amp; update firewall outbound egress filter rules.</li>
          </ul>
        </div>
      </div>
    </section>

    <!-- SLIDE 13 -->
    <section class="slide" data-slide="13">
      <div class="slide-badge purple">CASE STUDY 3 // LAYER 2 CYBER ATTACK</div>
      <h1 class="slide-title">The ARP Spoofing Doppelgänger</h1>
      <p class="slide-subtitle">Symptom: Multiple lab users report browser SSL/TLS certificate warnings &amp; slow speeds</p>

      <div class="grid-2">
        <div class="card">
          <h2 class="card-title">Inspecting the ARP Table</h2>
          <p class="card-desc">
            You run <code>arp -a</code> on a victim workstation and make an alarming discovery:
          </p>
          <div class="cli-box purple">Interface: 192.168.1.105 --- 0x3
  Internet Address      Physical Address      Type
  192.168.1.1           b8-27-eb-93-11-aa     dynamic  <-- Default Gateway
  192.168.1.55          b8-27-eb-93-11-aa     dynamic  <-- Student Laptop!</div>
          <p class="card-desc">
            <strong>Diagnosis:</strong> A rogue device (Raspberry Pi/Laptop) is sending gratuitous ARP replies, poisoning the LAN cache to conduct a <strong>Man-in-the-Middle (MitM) SSL stripping attack</strong>.
          </p>
        </div>

        <div class="card">
          <h2 class="card-title">Remediation &amp; Permanent Prevention</h2>
          <ul style="margin: 0.6rem 0 0 1.2rem; color: var(--text-muted); font-size: 0.85rem; line-height: 1.8;">
            <li><strong>Immediate Fix:</strong> Shut down switch port hosting MAC <code>b8-27-eb-93-11-aa</code>.</li>
            <li><strong>Client Remediation:</strong> Flush ARP cache with <code>netsh interface ip delete arpcache</code>.</li>
            <li><strong>Step 5 Prevention (Enterprise Hardening):</strong> Enable <strong>Dynamic ARP Inspection (DAI)</strong> and <strong>DHCP Snooping</strong> on Cisco switches so only authorized DHCP-assigned MACs can send ARP packets.</li>
            <li><strong>Step 6 Documentation:</strong> Document rogue MAC and report to school administration.</li>
          </ul>
        </div>
      </div>
    </section>

    <!-- SLIDE 14 -->
    <section class="slide" data-slide="14">
      <div class="slide-badge cyan">ENGINEERING STANDARDS</div>
      <h1 class="slide-title">How to Draw an Elite Logic Flowchart</h1>
      <p class="slide-subtitle">Turning your troubleshooting methodology into a bulletproof decision tree</p>

      <div class="grid-3">
        <div class="card">
          <div class="card-num cyan">01</div>
          <h2 class="card-title">Decision Diamonds</h2>
          <p class="card-desc">
            Every test in your flowchart must be a <strong>testable command</strong> with exactly two exit branches: <strong>[PASS &rarr; TRUE]</strong> and <strong>[FAIL &rarr; FALSE]</strong>.
          </p>
          <div class="cli-box cyan">[ping 192.168.1.1]
   /         \\
 [PASS]     [FAIL]</div>
        </div>

        <div class="card">
          <div class="card-num amber">02</div>
          <h2 class="card-title">No Ambiguous Steps</h2>
          <p class="card-desc">
            Never write <em>"Check computer"</em> or <em>"Fix internet"</em>. Write the exact CLI command: <code>ipconfig /flushdns</code>, <code>tracert -d 8.8.8.8</code>, or <code>net start spooler</code>.
          </p>
          <div class="card-tag" style="color:var(--amber);">Precision Engineering</div>
        </div>

        <div class="card">
          <div class="card-num">03</div>
          <h2 class="card-title">Terminal Endpoints</h2>
          <p class="card-desc">
            Every path must terminate in either:
            <br>&bull; <strong>[RESOLVED &amp; DOCUMENTED]</strong>
            <br>&bull; <strong>[ESCALATE TO TIER 2 / NET ADMIN]</strong>
          </p>
          <div class="card-tag">Definitive Closure</div>
        </div>
      </div>
    </section>

    <!-- SLIDE 15 -->
    <section class="slide" data-slide="15">
      <div class="slide-badge">MISSION BRIEFING</div>
      <h1 class="slide-title">Deploying to the Interactive Lab</h1>
      <p class="slide-subtitle">Ready your terminal. Open your companion activity worksheet.</p>

      <div class="grid-2">
        <div class="card" style="border: 2px solid var(--green);">
          <h2 class="card-title" style="color:var(--green); font-size:1.3rem;">Today's Lab Directives</h2>
          <ul style="margin: 0.8rem 0 0 1.2rem; color: var(--text-muted); font-size: 0.9rem; line-height: 2;">
            <li>1. Open the <strong>Divide &amp; Conquer Lab Activity</strong>.</li>
            <li>2. Work through the <strong>5 Progressive Scenarios</strong> (Physical &rarr; APIPA &rarr; DNS &rarr; Security Incident &rarr; ARP Attack).</li>
            <li>3. Use the built-in <strong>Virtual Terminal</strong> to test CLI commands.</li>
            <li>4. Construct your <strong>Logic Flowchart Trees</strong>.</li>
            <li>5. Export your completed work as a PDF for grading!</li>
          </ul>
        </div>

        <div class="card" style="justify-content:center; align-items:center; text-align:center; background:rgba(57,255,158,0.03);">
          <div style="font-size:3rem; margin-bottom:0.8rem;">⚡</div>
          <h3 style="font-size:1.3rem; color:#fff; font-family:'Space Grotesk';">Ready to Begin?</h3>
          <p style="color:var(--text-muted); font-size:0.85rem; margin:0.6rem 0 1.2rem;">Launch the interactive student workbench now.</p>
          <a href="activity.html" class="btn-ctrl" style="background:var(--green); color:#000; font-weight:800; font-size:0.9rem; padding:0.7rem 1.6rem; border-radius:6px; text-decoration:none;">
            [ OPEN LAB ACTIVITY HUB ]
          </a>
        </div>
      </div>
    </section>

  </main>

  <aside id="presenter-drawer" class="presenter-drawer">
    <div class="presenter-header">
      <span>TEACHER PRESENTER NOTES // MR. L'S GUIDE</span>
      <button onclick="toggleNotes()" style="background:none; border:none; color:var(--amber); cursor:pointer; font-weight:bold;">[CLOSE X]</button>
    </div>
    <div id="presenter-content" class="presenter-text">
      Welcome to Slide 1! Emphasize to students that troubleshooting is a structured engineering process, not magical guessing.
    </div>
  </aside>

  <footer class="bottombar">
    <div class="key-hints">
      <span><kbd>&larr;</kbd> Previous</span>
      <span><kbd>&rarr;</kbd> Next</span>
      <span><kbd>Space</kbd> Advance</span>
      <span><kbd>N</kbd> Notes</span>
      <span><kbd>F</kbd> Fullscreen</span>
    </div>

    <div class="progress-track">
      <div id="progress-fill" class="progress-fill"></div>
    </div>

    <div style="display:flex; align-items:center; gap:0.6rem;">
      <button class="btn-ctrl" onclick="prevSlide()" title="Previous Slide">&larr; PREV</button>
      <span id="slide-counter" class="slide-counter">01 / 15</span>
      <button class="btn-ctrl" onclick="nextSlide()" title="Next Slide">NEXT &rarr;</button>
    </div>
  </footer>

  <script>
    const PRESENTER_NOTES = {
      1: "SLIDE 1: Hook the class immediately. Ask students: 'When your Wi-Fi drops at home, what is the first thing you do?' Point out that pulling random plugs without knowing what layer failed is amateur. Introduce the 6 steps as the universal CompTIA standard.",
      2: "SLIDE 2: Business context. Emphasize downtime cost. In modern hospitals, data centers, and trading firms, network downtime is an existential crisis. Ask students what happens if you reboot a core switch during school testing.",
      3: "SLIDE 3: The Complete Lifecycle. Have the class repeat the 6 steps aloud: Identify, Theorize, Test, Plan/Implement, Verify/Prevent, Document. Highlight that Step 5 (preventative) and Step 6 (documentation) are where good techs become senior engineers.",
      4: "SLIDE 4: Step 1 (Identify). Teach students how to ask open-ended vs closed-ended questions. Ask: 'What is the danger of trusting the user when they say the server is down?' Check logs and verify scope first.",
      5: "SLIDE 5: Steps 2 & 3. Falsification! Explain that a failed test is great news because it permanently eliminates a wrong direction. Emphasize: DO NOT fix anything in Step 3; Step 3 is only for testing!",
      6: "SLIDE 6: Step 4 (Plan & Implement). Talk about Change Management in real enterprises. Mention outage windows and rollback procedures (e.g. taking a VM snapshot before installing an update).",
      7: "SLIDE 7: Steps 5 & 6. Verification must involve the user. Then explain preventative maintenance: why did it happen? If a cable failed, replace it with molded boot Cat6. Show ticket documentation format.",
      8: "SLIDE 8: Divide and Conquer Core Concept. Draw the binary search tree on the board. 7 linear tests take too long. By testing Layer 3, we cut the problem domain in half immediately.",
      9: "SLIDE 9: The 5-Step Midpoint Ladder. This is the Holy Grail! 1. 127.0.0.1 (NIC). 2. Gateway (LAN). 3. 8.8.8.8 (WAN). 4. google.com (DNS). 5. Port 443 (App). Have students copy this into their notebooks.",
      10: "SLIDE 10: Interactive Sandbox. Walk through the clicks live on the projector. Show what happens when Gateway passes vs fails.",
      11: "SLIDE 11: Case Study 1 (APIPA). Teach students why 169.254.0.0/16 happens. Link lights are green because physical Layer 1/2 is UP, but DHCP failed at Layer 3/7. Have them practice `ipconfig /renew`.",
      12: "SLIDE 12: Case Study 2 (Cybersecurity Incident). Transition into Ethical Hacking. Explain Port 22 SSH exfiltration. Immediate action: isolate endpoint from switch port before attacker can move laterally.",
      13: "SLIDE 13: Case Study 3 (ARP Poisoning). Look at `arp -a`. Two IPs with the same MAC address means a Man-in-the-Middle attacker! Introduce Dynamic ARP Inspection (DAI) and DHCP Snooping.",
      14: "SLIDE 14: Logic Flowcharts. Explain the rubric for today's lab activity. Every diamond must be an exact CLI command, and every branch must lead to pass/fail next steps.",
      15: "SLIDE 15: Mission Launch! Instruct students to open `activity.html` and begin Tier 1 through Tier 5 scenarios."
    };

    let currentSlide = 1;
    const totalSlides = 15;
    const slides = document.querySelectorAll('.slide');
    const counter = document.getElementById('slide-counter');
    const progressFill = document.getElementById('progress-fill');
    const presenterContent = document.getElementById('presenter-content');
    const presenterDrawer = document.getElementById('presenter-drawer');

    function updateSlide() {
      slides.forEach((s, idx) => {
        s.classList.remove('active', 'prev');
        if (idx + 1 === currentSlide) {
          s.classList.add('active');
        } else if (idx + 1 < currentSlide) {
          s.classList.add('prev');
        }
      });

      counter.textContent = `${String(currentSlide).padStart(2, '0')} / ${String(totalSlides).padStart(2, '0')}`;
      progressFill.style.width = `${(currentSlide / totalSlides) * 100}%`;

      if (PRESENTER_NOTES[currentSlide]) {
        presenterContent.textContent = PRESENTER_NOTES[currentSlide];
      }
    }

    function nextSlide() {
      if (currentSlide < totalSlides) {
        currentSlide++;
        updateSlide();
      }
    }

    function prevSlide() {
      if (currentSlide > 1) {
        currentSlide--;
        updateSlide();
      }
    }

    function toggleNotes() {
      presenterDrawer.classList.toggle('open');
      document.getElementById('btn-notes').classList.toggle('active');
    }

    function toggleFullScreen() {
      if (!document.fullscreenElement) {
        document.documentElement.requestFullscreen().catch(err => alert(err.message));
      } else {
        document.exitFullscreen();
      }
    }

    window.addEventListener('keydown', (e) => {
      if (e.key === 'ArrowRight' || e.key === 'Space' || e.key === 'PageDown') {
        e.preventDefault();
        nextSlide();
      } else if (e.key === 'ArrowLeft' || e.key === 'PageUp') {
        e.preventDefault();
        prevSlide();
      } else if (e.key.toLowerCase() === 'n') {
        toggleNotes();
      } else if (e.key.toLowerCase() === 'f') {
        toggleFullScreen();
      } else if (e.key === 'Home') {
        currentSlide = 1;
        updateSlide();
      } else if (e.key === 'End') {
        currentSlide = totalSlides;
        updateSlide();
      }
    });

    function runSim(step) {
      const output = document.getElementById('sim-output');
      const verdict = document.getElementById('sim-verdict');
      const s1 = document.getElementById('sim-status-1');
      const s2 = document.getElementById('sim-status-2');
      const s3 = document.getElementById('sim-status-3');
      const s4 = document.getElementById('sim-status-4');

      if (step === 1) {
        s1.className = 'osi-status pass'; s1.textContent = 'PASS';
        output.textContent = `C:\\\\> ping 127.0.0.1\\n\\nPinging 127.0.0.1 with 32 bytes of data:\\nReply from 127.0.0.1: bytes=32 time<1ms TTL=128\\nReply from 127.0.0.1: bytes=32 time<1ms TTL=128\\n\\nPing statistics for 127.0.0.1:\\n    Packets: Sent = 2, Received = 2, Lost = 0 (0% loss)`;
        verdict.textContent = "[+] LOOPBACK OK: Local TCP/IP software stack and NIC driver are 100% functional.";
      } else if (step === 2) {
        s2.className = 'osi-status pass'; s2.textContent = 'PASS';
        output.textContent = `C:\\\\> ping 192.168.1.1\\n\\nPinging 192.168.1.1 with 32 bytes of data:\\nReply from 192.168.1.1: bytes=32 time=1ms TTL=64\\nReply from 192.168.1.1: bytes=32 time=1ms TTL=64\\n\\n[+] LOCAL LAN LINK UP! Physical cable, wall jack, switch port, and router interface verified.`;
        verdict.textContent = "[+] GATEWAY REACHABLE: The problem is NOT your cable or switch! Problem space reduced by 50%.";
      } else if (step === 3) {
        s3.className = 'osi-status pass'; s3.textContent = 'PASS';
        output.textContent = `C:\\\\> ping 8.8.8.8\\n\\nPinging 8.8.8.8 with 32 bytes of data:\\nReply from 8.8.8.8: bytes=32 time=14ms TTL=117\\nReply from 8.8.8.8: bytes=32 time=12ms TTL=117\\n\\n[+] WAN IP ROUTING OK: Default route, NAT, and ISP modem connection are active!`;
        verdict.textContent = "[+] WAN REACHABLE: You have full external internet connectivity. Problem must be in Layer 7 (DNS/App).";
      } else if (step === 4) {
        s4.className = 'osi-status fail'; s4.textContent = 'FAIL';
        output.textContent = `C:\\\\> nslookup google.com\\n\\nServer:  UnKnown\\nAddress:  192.168.1.1\\n\\n*** DNS request to 192.168.1.1 timed out.\\n*** Request to UnKnown timed-out after 2 seconds.\\n\\n[-] ROOT CAUSE PINPOINTED: DNS Resolver timed out! Local IP routing is healthy, but DNS configuration is dead.`;
        verdict.innerHTML = "<span style='color:var(--red);'>[-] ROOT CAUSE FOUND:</span> Stale/Incorrect DNS Resolver IP in network adapter settings!";
      }
    }

    function setLanguage(lang) {
      document.documentElement.lang = lang;
      document.documentElement.dir = lang === 'ar' ? 'rtl' : 'ltr';
      localStorage.setItem('presentation_lang', lang);

      document.querySelectorAll('#lang-en, #lang-ar, #lang-uk').forEach(b => b.classList.remove('active'));
      const activeBtn = document.getElementById(`lang-${lang}`);
      if (activeBtn) activeBtn.classList.add('active');
    }

    updateSlide();
    const savedLang = localStorage.getItem('presentation_lang') || 'en';
    setLanguage(savedLang);
  </script>
</body>
</html>
"""

with open(presentation_html_path, "w", encoding="utf-8") as f:
    f.write(presentation_html)

print(f"Generated presentation.html -> {presentation_html_path}")

# ==========================================
# 2. GENERATE ACTIVITY.HTML
# ==========================================
activity_html_path = os.path.join(TARGET_DIR, "activity.html")

activity_html = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Lab Activity // CompTIA 6-Step Troubleshooting Workbench</title>
  
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=IBM+Plex+Sans+Arabic:wght@400;600;700&family=Inter:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;500;700;800&family=Space+Grotesk:wght@600;700&display=swap" rel="stylesheet">
  
  <style>
    :root {
      --bg: #040608;
      --panel: #090e13;
      --panel-elevated: #0f1720;
      --line: #1c2b26;
      --line-bright: #2d424b;
      --green: #39ff9e;
      --green-glow: rgba(57, 255, 158, 0.3);
      --green-dim: #164e37;
      --cyan: #38bdf8;
      --amber: #ffb454;
      --red: #ff5c72;
      --purple: #a855f7;
      --text: #d9f5e6;
      --text-muted: #79998d;
      --text-dim: #496359;
    }

    * { box-sizing: border-box; }

    html, body {
      margin: 0;
      padding: 0;
      min-height: 100%;
      background: var(--bg);
      color: var(--text);
      font-family: 'JetBrains Mono', 'Courier New', monospace;
    }

    :root:lang(ar) {
      font-family: 'IBM Plex Sans Arabic', 'Inter', sans-serif;
    }

    [dir="rtl"] { text-align: start; }
    [dir="rtl"] code, [dir="rtl"] pre, [dir="rtl"] .font-mono, [dir="rtl"] .cli-term, [dir="rtl"] input.cli-input {
      direction: ltr !important;
      unicode-bidi: isolate;
      text-align: left;
    }

    @media screen {
      body {
        background-image:
          linear-gradient(rgba(57,255,158,0.04) 1px, transparent 1px),
          linear-gradient(90deg, rgba(57,255,158,0.04) 1px, transparent 1px);
        background-size: 36px 36px;
        background-color: var(--bg);
        position: relative;
        overflow-x: hidden;
      }
      body::before {
        content: ""; position: fixed; inset: 0; pointer-events: none;
        background: repeating-linear-gradient(to bottom, rgba(0,0,0,0) 0px, rgba(0,0,0,0) 2px, rgba(0,0,0,0.18) 3px, rgba(0,0,0,0) 4px);
        mix-blend-mode: multiply; z-index: 5;
      }
      body::after {
        content: ""; position: fixed; inset: 0; pointer-events: none;
        background: radial-gradient(ellipse at center, transparent 45%, rgba(0,0,0,0.65) 100%);
        z-index: 4;
      }
      .panel {
        background: linear-gradient(180deg, var(--panel), #060a0d);
        border: 1px solid var(--line);
        border-radius: 8px;
        position: relative;
        box-shadow: 0 10px 30px rgba(0,0,0,0.7);
      }
      .panel::before {
        content: ""; position: absolute; top: -1px; left: -1px; right: -1px; bottom: -1px; border: 1px solid transparent;
        background: linear-gradient(120deg, rgba(57,255,158,0.35), transparent 30%, transparent 70%, rgba(56,189,248,0.3)) border-box;
        -webkit-mask: linear-gradient(#fff 0 0) padding-box, linear-gradient(#fff 0 0);
        -webkit-mask-composite: xor; mask-composite: exclude; pointer-events: none;
        border-radius: 8px;
      }
      .panel-head { border-bottom: 1px solid var(--line); }
      h1 {
        color: #eafff2;
        text-shadow: 0 0 8px rgba(57,255,158,0.35), 2px 0 0 rgba(255,92,114,0.25), -2px 0 0 rgba(84,180,255,0.25);
      }
      .student-input {
        background: rgba(4, 7, 10, 0.7);
        border: 1px solid var(--line);
        color: #fff;
        padding: 0.75rem 1rem;
        border-radius: 4px;
        font-family: inherit;
        font-size: 0.9rem;
        width: 100%;
        transition: all 0.2s ease;
      }
      .student-input:focus {
        outline: none;
        border-color: var(--green);
        box-shadow: 0 0 10px var(--green-glow);
        background: rgba(9, 14, 19, 0.95);
      }
    }

    .wrap { position: relative; z-index: 2; min-height: 100vh; display: flex; flex-direction: column; padding: 3vh 5vw 6vh; }
    
    .topbar {
      display: flex; justify-content: space-between; align-items: baseline;
      font-size: clamp(0.68rem, 1vw, 0.85rem); color: var(--text-muted); letter-spacing: 0.08em;
      border-bottom: 1px solid var(--line); padding-bottom: 0.8em; margin-bottom: 3vh;
    }
    .topbar a { color: var(--text-muted); text-decoration: none; }
    .topbar a:hover { color: var(--green); }
    .topbar .dot { color: var(--green); }
    .blink { animation: blink 1.1s steps(1) infinite; }
    @keyframes blink { 50% { opacity: 0; } }

    header { text-align: center; margin-bottom: 4vh; }
    .prompt-line { font-size: clamp(0.75rem, 1.1vw, 0.95rem); color: var(--green); margin-bottom: 0.8em; letter-spacing: 0.03em; }
    .prompt-line .path { color: var(--amber); }
    h1 {
      font-family: 'Space Grotesk', sans-serif;
      font-size: clamp(2rem, 4vw, 3.2rem); line-height: 1.1; margin: 0; font-weight: 700;
    }
    .subtitle { margin-top: 0.6em; margin-bottom: 1.5em; font-size: clamp(0.85rem, 1.2vw, 1.05rem); color: var(--text-muted); letter-spacing: 0.15em; text-transform: uppercase; }
    .subtitle::before, .subtitle::after { content: "//"; color: var(--green-dim); margin: 0 0.6em; }

    .action-bar {
      display: flex;
      justify-content: center;
      flex-wrap: wrap;
      gap: 1rem;
      margin-bottom: 2.5rem;
    }

    .btn-act {
      background: rgba(57,255,158,0.06);
      border: 1px solid var(--green);
      color: var(--green);
      font-family: 'JetBrains Mono', monospace;
      font-size: 0.82rem;
      font-weight: 700;
      letter-spacing: 0.08em;
      padding: 0.7em 1.4em;
      cursor: pointer;
      border-radius: 4px;
      transition: all 0.2s ease;
      display: inline-flex;
      align-items: center;
      gap: 0.5rem;
      text-decoration: none;
    }
    .btn-act:hover {
      background: rgba(57,255,158,0.18);
      box-shadow: 0 0 15px var(--green-glow);
      transform: translateY(-1px);
    }
    .btn-act.cyan { border-color: var(--cyan); color: var(--cyan); background: rgba(56,189,248,0.06); }
    .btn-act.cyan:hover { background: rgba(56,189,248,0.18); box-shadow: 0 0 15px rgba(56,189,248,0.4); }
    .btn-act.amber { border-color: var(--amber); color: var(--amber); background: rgba(255,180,84,0.06); }
    .btn-act.amber:hover { background: rgba(255,180,84,0.18); box-shadow: 0 0 15px rgba(255,180,84,0.4); }

    main { flex: 1; display: flex; flex-direction: column; align-items: center; gap: 2.5rem; }
    .panel { width: 100%; max-width: 1050px; }
    .panel-head {
      display: flex; align-items: center; justify-content: space-between; padding: 0.8em 1.4em;
      font-size: clamp(0.7rem, 1vw, 0.82rem); color: var(--text-muted); letter-spacing: 0.1em; text-transform: uppercase;
      font-weight: 700;
    }
    .panel-head .lights span { display: inline-block; width: 9px; height: 9px; border-radius: 50%; margin-right: 6px; }
    .lights .r { background: var(--red); } .lights .a { background: var(--amber); } .lights .g { background: var(--green); }

    .panel-body { padding: 2em; }
    
    .name-block {
      display: grid;
      grid-template-columns: 2fr 1fr 1fr;
      gap: 1.5rem;
      margin-bottom: 1.8rem;
      font-family: 'Space Grotesk', sans-serif;
    }
    .name-field { display: flex; flex-direction: column; gap: 0.3rem; }
    .name-field label { color: var(--green); font-weight: 700; font-family: 'JetBrains Mono'; font-size: 0.8rem; letter-spacing: 0.05em; }

    .scenario-title { font-family: 'Space Grotesk', sans-serif; font-size: 1.35rem; font-weight: 700; color: #fff; margin: 0 0 0.6em 0; display: flex; align-items: center; justify-content: space-between; }
    .scenario-title .tag { 
      font-size: 0.55em; vertical-align: middle; padding: 0.3em 0.7em; 
      border: 1px solid var(--amber); color: var(--amber); border-radius: 3px; font-family: 'JetBrains Mono'; font-weight: 800; letter-spacing: 0.1em;
    }
    .tag.hard { border-color: var(--red); color: var(--red); }
    .tag.easy { border-color: var(--green); color: var(--green); }
    .tag.cyan { border-color: var(--cyan); color: var(--cyan); }
    
    .situation {
      font-size: 0.95rem; color: var(--text); line-height: 1.6; margin-bottom: 1.5em; padding: 1rem 1.2rem;
      background: rgba(0,0,0,0.5); border-left: 3px solid var(--cyan); border-radius: 0 6px 6px 0;
    }
    
    .terminal-sim {
      background: #020305;
      border: 1px solid var(--line);
      border-radius: 6px;
      padding: 1rem 1.2rem;
      margin-bottom: 1.5rem;
    }
    .term-header {
      display: flex;
      justify-content: space-between;
      align-items: center;
      color: var(--text-muted);
      font-size: 0.75rem;
      margin-bottom: 0.8rem;
      border-bottom: 1px solid rgba(255,255,255,0.06);
      padding-bottom: 0.4rem;
    }
    .term-btn-group { display: flex; gap: 0.5rem; flex-wrap: wrap; margin-bottom: 0.8rem; }
    .term-btn {
      background: var(--panel-elevated);
      border: 1px solid var(--line-bright);
      color: var(--cyan);
      font-family: 'JetBrains Mono', monospace;
      font-size: 0.75rem;
      padding: 0.3rem 0.7rem;
      border-radius: 4px;
      cursor: pointer;
      transition: all 0.15s;
    }
    .term-btn:hover { border-color: var(--green); color: var(--green); }
    .term-output {
      background: #000;
      color: var(--green);
      padding: 0.8rem;
      border-radius: 4px;
      font-size: 0.82rem;
      min-height: 80px;
      max-height: 180px;
      overflow-y: auto;
      white-space: pre-wrap;
    }

    .step-block { margin-bottom: 1.8rem; }
    .step-label {
      font-size: 0.85rem; font-weight: 800; color: var(--green); margin-bottom: 0.4em; letter-spacing: 0.06em;
      display: flex; align-items: center; gap: 0.5rem;
    }
    .step-label.cyan { color: var(--cyan); }
    .step-label.amber { color: var(--amber); }
    .step-prompt { font-size: 0.9rem; color: var(--text-muted); margin-bottom: 0.6em; line-height: 1.5; }

    /* Teacher Key Callout */
    .teacher-key-box {
      display: none;
      background: rgba(255, 180, 84, 0.08);
      border: 1px solid var(--amber);
      border-radius: 6px;
      padding: 1.2rem 1.4rem;
      margin-top: 1.5rem;
      font-size: 0.85rem;
      color: #fed7aa;
      line-height: 1.6;
    }
    .teacher-key-box.visible { display: block; }
    .teacher-key-title {
      color: var(--amber);
      font-weight: 800;
      letter-spacing: 0.08em;
      margin-bottom: 0.5rem;
      display: flex;
      align-items: center;
      gap: 0.5rem;
      text-transform: uppercase;
      font-size: 0.8rem;
    }

    footer { text-align: center; margin-top: 4vh; color: var(--text-dim); font-size: clamp(0.72rem, 1vw, 0.86rem); letter-spacing: 0.06em; }
    footer .cursor { color: var(--green); }

    /* PRINTER OPTIMIZED CSS */
    @media print {
      body { background: #ffffff !important; color: #000000 !important; font-family: Arial, sans-serif !important; }
      .topbar, footer, .prompt-line, header h1, .subtitle, .action-bar, .term-btn-group, .term-output { display: none !important; }
      .wrap { padding: 0 !important; }
      main { gap: 1.5rem !important; }
      .panel { background: #ffffff !important; border: 1px solid #000000 !important; page-break-inside: avoid; margin-bottom: 25px; box-shadow: none !important; }
      .panel::before { display: none !important; }
      .panel-head { border-bottom: 1px solid #000000 !important; color: #000000 !important; background: #f0f0f0 !important; padding: 0.5em 1em !important; }
      .lights { display: none !important; }
      .scenario-title { color: #000000 !important; font-size: 1.2rem !important; font-weight: bold !important; }
      .situation { color: #222222 !important; border-left: 3px solid #666666 !important; background: #fafafa !important; padding: 0.8em !important; }
      .step-label { color: #000000 !important; font-weight: bold !important; font-size: 0.85rem !important; }
      .step-prompt { color: #333333 !important; font-size: 0.85rem !important; }
      .student-input {
        background: transparent !important;
        border: none !important;
        border-bottom: 1px solid #777777 !important;
        border-radius: 0 !important;
        color: #000000 !important;
        padding: 0.2rem 0 !important;
        font-size: 0.85rem !important;
        min-height: 2em !important;
      }
      .tag { border-color: #000000 !important; color: #000000 !important; }
      .print-header { display: block !important; text-align: center; margin-bottom: 1.5rem; border-bottom: 2px solid #000; padding-bottom: 0.5rem; }
      .print-header h2 { margin: 0; font-size: 1.4rem; }
      .print-header p { margin: 0.2rem 0 0; font-size: 0.85rem; color: #444; }
    }

    .print-header { display: none; }
  </style>
</head>
<body>
  <div class="wrap">
    
    <!-- Printable Header (Print Only) -->
    <div class="print-header">
      <h2>A.W. Beattie Career Center &middot; Network Engineering &amp; Cybersecurity</h2>
      <p>Lab Activity: Divide &amp; Conquer // The CompTIA 6-Step Troubleshooting Model</p>
    </div>

    <div class="topbar">
      <span><span class="dot">●</span> BEATTIE-NET // NETWORK ENGINEERING &amp; CYBERSECURITY</span>
      <span>MR. L <span class="blink">_</span></span>
    </div>

    <header>
      <div class="prompt-line"><span class="path">mrl@beattie-tech</span>:~$ ./run_activity.sh --module="Troubleshooting"</div>
      <h1>Divide &amp; Conquer Lab Workbench</h1>
      <div class="subtitle">The CompTIA 6-Step Troubleshooting Model</div>
      
      <div class="action-bar">
        <button class="btn-act" onclick="window.print()">
          📄 <span>Export / Print Student Worksheet (PDF)</span>
        </button>

        <button id="btn-toggle-key" class="btn-act amber" onclick="toggleAnswerKey()">
          🔑 <span id="key-label">Show Teacher Answer Key &amp; Rubric</span>
        </button>

        <a href="presentation.html" class="btn-act cyan">
          📺 <span>Open Slide Presentation &rarr;</span>
        </a>

        <!-- 3-Way Language Toggle -->
        <div style="display:flex; background:var(--panel-elevated); border:1px solid var(--line-bright); border-radius:4px; padding:2px; font-size:0.75rem;">
          <button id="lang-en" class="btn-act" onclick="setLanguage('en')" style="border:none; padding:0.3rem 0.6rem;">EN</button>
          <button id="lang-ar" class="btn-act" onclick="setLanguage('ar')" style="border:none; padding:0.3rem 0.6rem;">عربي</button>
          <button id="lang-uk" class="btn-act" onclick="setLanguage('uk')" style="border:none; padding:0.3rem 0.6rem;">УКР</button>
        </div>
      </div>
    </header>

    <main>
      
      <!-- INSTRUCTIONS & STUDENT IDENTIFICATION -->
      <div class="panel">
        <div class="panel-head">
          <span class="lights"><span class="r"></span><span class="a"></span><span class="g"></span> auth.log</span>
          <span>STUDENT IDENTIFICATION &amp; METHODOLOGY RUBRIC</span>
        </div>
        <div class="panel-body">
          <div class="name-block">
            <div class="name-field">
              <label>STUDENT NAME</label>
              <input id="stud-name" type="text" class="student-input" placeholder="e.g. John Doe">
            </div>
            <div class="name-field">
              <label>DATE</label>
              <input id="stud-date" type="date" class="student-input">
            </div>
            <div class="name-field">
              <label>PERIOD / SECTION</label>
              <input id="stud-period" type="text" class="student-input" placeholder="e.g. AM-1">
            </div>
          </div>

          <div style="font-size:0.88rem; color:var(--text-muted); border: 1px dashed var(--line); padding: 1.2rem; border-radius: 6px; line-height: 1.7; background: rgba(0,0,0,0.4);">
            <strong style="color:var(--green);">MISSION DIRECTIVE:</strong> For each scenario below, apply the formal <strong>CompTIA 6-Step Troubleshooting Model</strong> using the <strong>Divide &amp; Conquer (Binary Search)</strong> methodology. Use the interactive virtual terminals to run diagnostic commands, then document your findings and action plan in the structured fields below.
            <div style="margin-top: 0.6rem; color: var(--cyan); font-size: 0.8rem;">
              [1] Identify problem &rarr; [2] Establish theory &rarr; [3] Test theory &rarr; [4] Plan &amp; implement fix &rarr; [5] Verify &amp; prevent &rarr; [6] Document outcomes.
            </div>
          </div>
        </div>
      </div>

      <!-- SCENARIO 1 -->
      <div class="panel" id="sec-scen-1">
        <div class="panel-head">
          <span class="lights"><span class="r"></span><span class="a"></span><span class="g"></span> scenario_01.log</span>
          <span>LEVEL 1: PHYSICAL LAYER // LAYER 1 &amp; 2</span>
        </div>
        <div class="panel-body">
          <div class="scenario-title">
            <span>The Silent Jack</span>
            <span class="tag easy">DIFFICULTY: NOVICE &middot; 100 PTS</span>
          </div>
          
          <div class="situation">
            <strong>SITUATION:</strong> A student punches down a new Cat6 UTP cable into patch panel port 14 and wires the keystone wall jack in Room 204. They plug a test workstation into the jack using a known-good patch cable. The NIC link LEDs stay completely dark, and Windows reports: <code>"Network Cable Unplugged."</code>
          </div>

          <!-- Virtual Terminal Sandbox -->
          <div class="terminal-sim">
            <div class="term-header">
              <span>VIRTUAL TERMINAL // WORKSTATION-ROOM204</span>
              <span>CLI DIAGNOSTIC SANDBOX</span>
            </div>
            <div class="term-btn-group">
              <button class="term-btn" onclick="runTerm(1, 'ipconfig')">run: ipconfig /all</button>
              <button class="term-btn" onclick="runTerm(1, 'link')">run: Test-NetConnection -InterfaceIndex 3</button>
              <button class="term-btn" onclick="runTerm(1, 'tester')">run: check_cable_tester_output</button>
            </div>
            <div id="term-out-1" class="term-output">Click a diagnostic button above to inspect endpoint telemetry...</div>
          </div>
          
          <div class="step-block">
            <div class="step-label">[STEP 1 &amp; 2] // IDENTIFY PROBLEM &amp; ESTABLISH THEORY</div>
            <div class="step-prompt">What physical components must be checked first? State two probable physical root causes for this Layer 1 failure.</div>
            <input id="q1-1" type="text" class="student-input" placeholder="Enter your identification and theory...">
          </div>
          
          <div class="step-block">
            <div class="step-label cyan">[STEP 3 &amp; 4] // TEST THEORY &amp; IMPLEMENT ACTION PLAN</div>
            <div class="step-prompt">What hardware tool will you use to test the cable run? If the wiremap shows pins 1 and 2 are reversed (miswire), what is your step-by-step fix?</div>
            <input id="q1-2" type="text" class="student-input" placeholder="Enter test tool and remediation steps...">
          </div>
          
          <div class="step-block">
            <div class="step-label amber">[STEP 5 &amp; 6] // VERIFY SYSTEM &amp; DOCUMENT IN TICKET</div>
            <div class="step-prompt">How do you verify full functionality (Step 5), and what exact resolution text should be recorded in the helpdesk ticket (Step 6)?</div>
            <input id="q1-3" type="text" class="student-input" placeholder="Enter verification test and ticket resolution log...">
          </div>

          <!-- Teacher Answer Key Box -->
          <div class="teacher-key-box" id="key-1">
            <div class="teacher-key-title">🔑 TEACHER ANSWER KEY &amp; RUBRIC // SCENARIO 1</div>
            <strong>Step 1 &amp; 2:</strong> Check NIC link lights, patch cable on both ends, switch port status. Theories: Keystone jack punchdown wire reversed (miswire), cable split pair, bad punchdown blade connection, or switch port disabled.<br>
            <strong>Step 3 &amp; 4:</strong> Use a digital Cable Continuity Tester / Wiremapper. If pins 1 &amp; 2 are crossed, re-punch the keystone jack following TIA/EIA-568B standard using a 110 punchdown tool with cut blade facing out.<br>
            <strong>Step 5 &amp; 6:</strong> Verify link light goes solid green at 1 Gbps, workstation acquires IP via DHCP, and pings gateway. Ticket Doc: <em>"Re-punched Room 204 Port 14 keystone jack to T-568B to resolve pin 1-2 cross. Verified 1Gbps link & DHCP acquisition."</em>
          </div>
        </div>
      </div>

      <!-- SCENARIO 2 -->
      <div class="panel" id="sec-scen-2">
        <div class="panel-head">
          <span class="lights"><span class="r"></span><span class="a"></span><span class="g"></span> scenario_02.log</span>
          <span>LEVEL 2: NETWORK LAYER // LAYER 3 (DHCP &amp; IPV4)</span>
        </div>
        <div class="panel-body">
          <div class="scenario-title">
            <span>The APIPA 169.254 Anomaly</span>
            <span class="tag easy">DIFFICULTY: NOVICE &middot; 150 PTS</span>
          </div>
          
          <div class="situation">
            <strong>SITUATION:</strong> All 12 computers in Lab B suddenly lose network connectivity simultaneously. Physical link LEDs on both the PCs and the switch are green. You run <code>ipconfig</code> on PC-01 and see IPv4 address <code>169.254.88.204</code> with subnet mask <code>255.255.0.0</code> and a blank default gateway.
          </div>

          <div class="terminal-sim">
            <div class="term-header">
              <span>VIRTUAL TERMINAL // LAB-B-PC01</span>
              <span>CLI DIAGNOSTIC SANDBOX</span>
            </div>
            <div class="term-btn-group">
              <button class="term-btn" onclick="runTerm(2, 'ipconfig')">run: ipconfig /all</button>
              <button class="term-btn" onclick="runTerm(2, 'ping_loop')">run: ping 127.0.0.1</button>
              <button class="term-btn" onclick="runTerm(2, 'renew')">run: ipconfig /renew</button>
              <button class="term-btn" onclick="runTerm(2, 'dhcp_status')">run: Get-DhcpServerv4Scope -ComputerName "DC01"</button>
            </div>
            <div id="term-out-2" class="term-output">Click a diagnostic button above to inspect endpoint telemetry...</div>
          </div>
          
          <div class="step-block">
            <div class="step-label">[STEP 1 &amp; 2] // IDENTIFY PROBLEM &amp; ESTABLISH THEORY</div>
            <div class="step-prompt">What does the 169.254.x.x address indicate? Since link lights are on, what network service is failing?</div>
            <input id="q2-1" type="text" class="student-input" placeholder="Explain APIPA and identify the failed service...">
          </div>
          
          <div class="step-block">
            <div class="step-label cyan">[STEP 3 &amp; 4] // TEST THEORY &amp; IMPLEMENT ACTION PLAN</div>
            <div class="step-prompt">How do you test the DHCP service? If the DHCP scope pool is 100% full, what is your remediation plan?</div>
            <input id="q2-2" type="text" class="student-input" placeholder="Explain testing command and how to fix scope exhaustion...">
          </div>
          
          <div class="step-block">
            <div class="step-label amber">[STEP 5 &amp; 6] // VERIFY SYSTEM &amp; PREVENTATIVE HARDENING</div>
            <div class="step-prompt">What CLI command forces the clients to get a new IP? What preventative measure stops this from recurring?</div>
            <input id="q2-3" type="text" class="student-input" placeholder="Command to force IP lease and preventative setting...">
          </div>

          <div class="teacher-key-box" id="key-2">
            <div class="teacher-key-title">🔑 TEACHER ANSWER KEY &amp; RUBRIC // SCENARIO 2</div>
            <strong>Step 1 &amp; 2:</strong> 169.254.x.x is APIPA (Automatic Private IP Addressing / RFC 3927). It proves Layer 1 & 2 are UP, but the client broadcasted DHCPDISCOVER and received no DHCPOFFER. The DHCP Server service is offline or scope is exhausted.<br>
            <strong>Step 3 &amp; 4:</strong> Check DHCP Server console/service. If scope is exhausted, expand the subnet (e.g. from /24 to /23) or reduce DHCP lease time from 8 days to 4 hours to reclaim stale leases.<br>
            <strong>Step 5 &amp; 6:</strong> Run <code>ipconfig /release</code> followed by <code>ipconfig /renew</code>. Preventative: configure automated 80% scope capacity threshold alerts in Windows Server / Syslog.
          </div>
        </div>
      </div>

      <!-- SCENARIO 3 -->
      <div class="panel" id="sec-scen-3">
        <div class="panel-head">
          <span class="lights"><span class="r"></span><span class="a"></span><span class="g"></span> scenario_03.log</span>
          <span>LEVEL 3: TRANSPORT &amp; APPLICATION LAYER // DNS &amp; FIREWALL</span>
        </div>
        <div class="panel-body">
          <div class="scenario-title">
            <span>The DNS Dilemma</span>
            <span class="tag cyan">DIFFICULTY: INTERMEDIATE &middot; 200 PTS</span>
          </div>
          
          <div class="situation">
            <strong>SITUATION:</strong> A teacher reports: <em>"The internet is completely down in my classroom!"</em> You open the command prompt on her workstation. You can successfully <code>ping 8.8.8.8</code> with 12ms latency, but when you run <code>ping google.com</code>, you get: <code>"Ping request could not find host google.com. Please check the name and try again."</code>
          </div>

          <div class="terminal-sim">
            <div class="term-header">
              <span>VIRTUAL TERMINAL // TEACHER-WS01</span>
              <span>CLI DIAGNOSTIC SANDBOX</span>
            </div>
            <div class="term-btn-group">
              <button class="term-btn" onclick="runTerm(3, 'ping_ip')">run: ping 8.8.8.8</button>
              <button class="term-btn" onclick="runTerm(3, 'ping_dns')">run: ping google.com</button>
              <button class="term-btn" onclick="runTerm(3, 'nslookup')">run: nslookup google.com</button>
              <button class="term-btn" onclick="runTerm(3, 'ipconfig_dns')">run: ipconfig /displaydns</button>
            </div>
            <div id="term-out-3" class="term-output">Click a diagnostic button above to inspect endpoint telemetry...</div>
          </div>
          
          <div class="step-block">
            <div class="step-label">[STEP 1 &amp; 2] // IDENTIFY PROBLEM &amp; ESTABLISH THEORY</div>
            <div class="step-prompt">What does the successful 8.8.8.8 ping prove about Layers 1, 2, and 3? What protocol is failing on google.com?</div>
            <input id="q3-1" type="text" class="student-input" placeholder="Analyze what the successful vs failed ping proves...">
          </div>
          
          <div class="step-block">
            <div class="step-label cyan">[STEP 3 &amp; 4] // TEST THEORY &amp; IMPLEMENT ACTION PLAN</div>
            <div class="step-prompt">What command tests DNS resolution specifically? If the adapter had a rogue DNS IP (10.0.0.99) manually configured, how do you fix it?</div>
            <input id="q3-2" type="text" class="student-input" placeholder="Testing command and configuration fix...">
          </div>
          
          <div class="step-block">
            <div class="step-label amber">[STEP 5 &amp; 6] // VERIFY SYSTEM &amp; PREVENTATIVE HARDENING</div>
            <div class="step-prompt">What command clears the local DNS cache? Write the exact 1-sentence resolution note for the ticket.</div>
            <input id="q3-3" type="text" class="student-input" placeholder="Cache clear command and ticket resolution text...">
          </div>

          <div class="teacher-key-box" id="key-3">
            <div class="teacher-key-title">🔑 TEACHER ANSWER KEY &amp; RUBRIC // SCENARIO 3</div>
            <strong>Step 1 &amp; 2:</strong> Successful ping to 8.8.8.8 proves physical link, local gateway, NAT, and WAN internet routing (Layers 1-3) are 100% operational! The failure on hostname proves DNS (Domain Name System, UDP/TCP Port 53) at Layer 7 is failing.<br>
            <strong>Step 3 &amp; 4:</strong> Test with <code>nslookup google.com</code>. Fix: Set IPv4 adapter properties to <em>"Obtain DNS server address automatically"</em> via DHCP or set to verified enterprise DNS servers (e.g. 1.1.1.1, 8.8.8.8).<br>
            <strong>Step 5 &amp; 6:</strong> Run <code>ipconfig /flushdns</code>. Ticket Doc: <em>"Removed static rogue DNS address 10.0.0.99; set adapter to DHCP DNS auto-discovery and flushed local DNS resolver cache. Verified web browsing to google.com."</em>
          </div>
        </div>
      </div>

      <!-- SCENARIO 4 -->
      <div class="panel" id="sec-scen-4">
        <div class="panel-head">
          <span class="lights"><span class="r"></span><span class="a"></span><span class="g"></span> scenario_04.log</span>
          <span>LEVEL 4: INCIDENT RESPONSE // CYBERSECURITY TRIAGE</span>
        </div>
        <div class="panel-body">
          <div class="scenario-title">
            <span>Suspicious Outbound Port 22 Exfiltration</span>
            <span class="tag hard">DIFFICULTY: ADVANCED &middot; 250 PTS</span>
          </div>
          
          <div class="situation">
            <strong>SITUATION:</strong> Your Security Operations Center (SOC) IDS fires a critical severity alert. Workstation-19 (192.168.10.84) in the CAD lab has established an active outbound connection on <strong>Port 22 (SSH)</strong> to an unclassified IP in Bulgaria (185.220.101.5), transferring over 3.2 GB of compressed data.
          </div>

          <div class="terminal-sim">
            <div class="term-header">
              <span>VIRTUAL TERMINAL // WORKSTATION-19 FORENSIC SHELL</span>
              <span>CLI DIAGNOSTIC SANDBOX</span>
            </div>
            <div class="term-btn-group">
              <button class="term-btn" onclick="runTerm(4, 'netstat')">run: netstat -ano | findstr :22</button>
              <button class="term-btn" onclick="runTerm(4, 'tasklist')">run: tasklist /fi "PID eq 4820"</button>
              <button class="term-btn" onclick="runTerm(4, 'firewall')">run: netsh advfirewall firewall show rule name=all</button>
            </div>
            <div id="term-out-4" class="term-output">Click a diagnostic button above to inspect endpoint telemetry...</div>
          </div>
          
          <div class="step-block">
            <div class="step-label">[STEP 1 &amp; 2] // IDENTIFY PROBLEM &amp; ESTABLISH THEORY</div>
            <div class="step-prompt">What is your IMMEDIATE first operational action (Step 1)? What protocol runs on Port 22, and what is the attacker doing?</div>
            <input id="q4-1" type="text" class="student-input" placeholder="Immediate action and threat theory...">
          </div>
          
          <div class="step-block">
            <div class="step-label cyan">[STEP 3 &amp; 4] // TEST THEORY &amp; IMPLEMENT ACTION PLAN</div>
            <div class="step-prompt">What CLI command links the network socket to the rogue process PID? How do you terminate the process and neutralize persistence?</div>
            <input id="q4-2" type="text" class="student-input" placeholder="Process inspection command and neutralization steps...">
          </div>
          
          <div class="step-block">
            <div class="step-label amber">[STEP 5 &amp; 6] // VERIFY SYSTEM &amp; PREVENTATIVE HARDENING</div>
            <div class="step-prompt">What firewall egress rule should be implemented to prevent unauthorized outbound SSH from student workstations in the future?</div>
            <input id="q4-3" type="text" class="student-input" placeholder="Firewall egress policy and post-incident documentation...">
          </div>

          <div class="teacher-key-box" id="key-4">
            <div class="teacher-key-title">🔑 TEACHER ANSWER KEY &amp; RUBRIC // SCENARIO 4</div>
            <strong>Step 1 &amp; 2:</strong> Immediate Action: Isolate the endpoint from the network (unplug cable or shut switch port) to prevent further data loss or lateral movement. Port 22 is SSH (Secure Shell). Attacker is running a reverse SSH tunnel for data exfiltration.<br>
            <strong>Step 3 &amp; 4:</strong> Run <code>netstat -ano | findstr :22</code> to identify the PID (e.g. PID 4820), then <code>taskkill /PID 4820 /F</code>. Inspect Task Scheduler and Startup registry keys for persistence scripts.<br>
            <strong>Step 5 &amp; 6:</strong> Preventative: Block outbound Port 22 egress on edge firewall for all subnets except designated bastion hosts. Run malware scan, re-image system from golden image, and document incident response report.
          </div>
        </div>
      </div>

      <!-- SCENARIO 5 -->
      <div class="panel" id="sec-scen-5">
        <div class="panel-head">
          <span class="lights"><span class="r"></span><span class="a"></span><span class="g"></span> scenario_05.log</span>
          <span>LEVEL 5: ENTERPRISE SECURITY // LAYER 2 CYBER ATTACK</span>
        </div>
        <div class="panel-body">
          <div class="scenario-title">
            <span>The Doppelgänger ARP Attack</span>
            <span class="tag hard">DIFFICULTY: EXPERT &middot; 300 PTS</span>
          </div>
          
          <div class="situation">
            <strong>SITUATION:</strong> Users on VLAN 10 report random web browser SSL/TLS certificate warnings (e.g. <em>"Your connection is not private - invalid certificate authority"</em>) and dropped connections. You inspect the victim workstation and check the ARP cache.
          </div>

          <div class="terminal-sim">
            <div class="term-header">
              <span>VIRTUAL TERMINAL // VICTIM-WS04</span>
              <span>CLI DIAGNOSTIC SANDBOX</span>
            </div>
            <div class="term-btn-group">
              <button class="term-btn" onclick="runTerm(5, 'arp')">run: arp -a</button>
              <button class="term-btn" onclick="runTerm(5, 'tracert')">run: tracert -d 8.8.8.8</button>
              <button class="term-btn" onclick="runTerm(5, 'switch_mac')">run: show mac address-table | include b827</button>
            </div>
            <div id="term-out-5" class="term-output">Click a diagnostic button above to inspect endpoint telemetry...</div>
          </div>
          
          <div class="step-block">
            <div class="step-label">[STEP 1 &amp; 2] // IDENTIFY PROBLEM &amp; ESTABLISH THEORY</div>
            <div class="step-prompt">Look at the <code>arp -a</code> table. Why do two different IP addresses have the exact same MAC address? Name the exact Layer 2 cyber attack.</div>
            <input id="q5-1" type="text" class="student-input" placeholder="Explain the ARP anomaly and name the cyber attack...">
          </div>
          
          <div class="step-block">
            <div class="step-label cyan">[STEP 3 &amp; 4] // TEST THEORY &amp; IMPLEMENT ACTION PLAN</div>
            <div class="step-prompt">How do you trace the rogue MAC to a physical switch port? What immediate action stops the attack?</div>
            <input id="q5-2" type="text" class="student-input" placeholder="Switch tracing command and containment action...">
          </div>
          
          <div class="step-block">
            <div class="step-label amber">[STEP 5 &amp; 6] // VERIFY SYSTEM &amp; ENTERPRISE HARDENING</div>
            <div class="step-prompt">What two enterprise Cisco switch security features (Step 5) permanently block ARP spoofing and rogue DHCP servers forever?</div>
            <input id="q5-3" type="text" class="student-input" placeholder="Enterprise switch security features (DAI / DHCP Snooping)...">
          </div>

          <div class="teacher-key-box" id="key-5">
            <div class="teacher-key-title">🔑 TEACHER ANSWER KEY &amp; RUBRIC // SCENARIO 5</div>
            <strong>Step 1 &amp; 2:</strong> ARP Spoofing / ARP Cache Poisoning (Man-in-the-Middle). A rogue device is emitting gratuitous ARP replies claiming the Default Gateway IP (192.168.1.1) belongs to its own MAC address (b8-27-eb-93-11-aa) to intercept and decrypt user traffic.<br>
            <strong>Step 3 &amp; 4:</strong> Run <code>show mac address-table | include b827</code> on core switch to find the physical port (e.g. GigabitEthernet 0/12). Issue <code>shutdown</code> command on that port to sever the attacker.<br>
            <strong>Step 5 &amp; 6:</strong> Enable <strong>Dynamic ARP Inspection (DAI)</strong> and <strong>DHCP Snooping</strong> on all access switches, combined with <strong>Port Security (Sticky MAC)</strong>. Flush victim ARP cache with <code>netsh interface ip delete arpcache</code>. Document in security incident log.
          </div>
        </div>
      </div>

    </main>

    <footer>
      END OF TRANSMISSION <span class="cursor blink">█</span>
    </footer>
  </div>

  <script>
    // Terminal Simulations Database
    const TERM_DATA = {
      1: {
        'ipconfig': "Windows IP Configuration\\n\\nEthernet adapter Ethernet0:\\n   Media State . . . . . . . . . . . : Media disconnected\\n   Connection-specific DNS Suffix  . : \\n   Description . . . . . . . . . . . : Intel(R) I211 Gigabit Network Connection",
        'link': "Test-NetConnection -InterfaceIndex 3\\n\\nInterfaceAlias               : Ethernet0\\nInterfaceIndex               : 3\\nOperationalStatus            : Down\\nMediaConnectionState         : Disconnected",
        'tester': "HARDWARE CABLE TESTER REPORT (Klein Tools Scout Pro 3):\\n[PASS] Length: 42 meters (Cat6 UTP)\\n[FAIL] WIREMAP ANOMALY:\\n   Tester Pin 1 -> Remote Pin 2\\n   Tester Pin 2 -> Remote Pin 1\\n   Pins 3, 4, 5, 6, 7, 8 -> Straight Through (OK)\\nRESULT: MISWIRE (Reversed Pair 1-2 on Room 204 Keystone Jack)"
      },
      2: {
        'ipconfig': "Windows IP Configuration\\n\\nEthernet adapter Ethernet0:\\n   Connection-specific DNS Suffix  . : \\n   IPv4 Address. . . . . . . . . . . : 169.254.88.204\\n   Subnet Mask . . . . . . . . . . . : 255.255.0.0\\n   Default Gateway . . . . . . . . . : \\n   DHCP Enabled. . . . . . . . . . . : Yes\\n   Autoconfiguration Enabled . . . . : Yes",
        'ping_loop': "Pinging 127.0.0.1 with 32 bytes of data:\\nReply from 127.0.0.1: bytes=32 time<1ms TTL=128\\nReply from 127.0.0.1: bytes=32 time<1ms TTL=128\\n[+] TCP/IP Software Stack is HEALTHY.",
        'renew': "C:\\\\> ipconfig /renew\\n\\nAn error occurred while renewing interface Ethernet0 : unable to contact your DHCP server. Request has timed out.",
        'dhcp_status': "PS C:\\\\> Get-DhcpServerv4Scope -ComputerName 'DC01'\\n\\nScopeId         SubnetMask      Name           State    StartRange      EndRange        LeasesInUse  AvailableLeases\\n-------         ----------      ----           -----    ----------      --------        -----------  ---------------\\n192.168.10.0    255.255.255.0   Lab_B_VLAN     Active   192.168.10.50   192.168.10.254  204 (100%)   0 (0% EXHAUSTED)"
      },
      3: {
        'ping_ip': "Pinging 8.8.8.8 with 32 bytes of data:\\nReply from 8.8.8.8: bytes=32 time=12ms TTL=117\\nReply from 8.8.8.8: bytes=32 time=11ms TTL=117\\n[+] WAN IP Routing and Gateway are 100% OPERATIONAL.",
        'ping_dns': "C:\\\\> ping google.com\\nPing request could not find host google.com. Please check the name and try again.",
        'nslookup': "C:\\\\> nslookup google.com\\nServer:  UnKnown\\nAddress:  10.0.0.99\\n\\n*** DNS request to 10.0.0.99 timed out.\\n*** Request to UnKnown timed-out after 2 seconds.\\n*** DNS FAILURE: Host unreachable.",
        'ipconfig_dns': "Windows IP Configuration\\n\\n   DNS Servers . . . . . . . . . . . : 10.0.0.99 (STALE / INVALID STATIC OVERRIDE)\\n   NetBIOS over Tcpip. . . . . . . . : Enabled"
      },
      4: {
        'netstat': "Active Connections\\n  Proto  Local Address          Foreign Address        State           PID\\n  TCP    192.168.10.84:54210    185.220.101.5:22       ESTABLISHED     4820\\n  TCP    192.168.10.84:139      0.0.0.0:0              LISTENING       4",
        'tasklist': "Image Name                     PID Session Name        Session#    Mem Usage\\n========================= ======== ================ =========== ============\\nputty_persist.exe             4820 Console                    1     14,240 K",
        'firewall': "Rule Name: Rogue Outbound SSH Rule\\n----------------------------------------------------------------------\\nEnabled: Yes\\nDirection: Out\\nProfiles: Domain,Private,Public\\nAction: Allow\\nProtocol: TCP\\nRemotePort: 22"
      },
      5: {
        'arp': "Interface: 192.168.1.105 --- 0x3\\n  Internet Address      Physical Address      Type\\n  192.168.1.1           b8-27-eb-93-11-aa     dynamic   <-- (DEFAULT GATEWAY)\\n  192.168.1.55          b8-27-eb-93-11-aa     dynamic   <-- (STUDENT LAPTOP - DUPLICATE MAC!)\\n  192.168.1.100         00-1a-2b-3c-4d-5e     dynamic",
        'tracert': "Tracing route to 8.8.8.8 over a maximum of 30 hops:\\n  1    <1 ms    <1 ms    <1 ms  192.168.1.55 (MITM PROXY INTERCEPT!)\\n  2     1 ms     1 ms     1 ms  192.168.1.1\\n  3    12 ms    11 ms    11 ms  96.120.45.1",
        'switch_mac': "Core-Switch-01# show mac address-table | include b827\\nVLAN    Mac Address       Type       Ports\\n----    -----------       ----       -----\\n10      b827.eb93.11aa    DYNAMIC    Gi0/12  (Room 104 Drop 6)"
      }
    };

    function runTerm(scenNum, cmdKey) {
      const outputElem = document.getElementById(`term-out-${scenNum}`);
      if (TERM_DATA[scenNum] && TERM_DATA[scenNum][cmdKey]) {
        outputElem.textContent = TERM_DATA[scenNum][cmdKey];
      }
    }

    function toggleAnswerKey() {
      const keyBoxes = document.querySelectorAll('.teacher-key-box');
      const isVisible = keyBoxes[0].classList.contains('visible');
      keyBoxes.forEach(b => b.classList.toggle('visible', !isVisible));
      document.getElementById('key-label').textContent = isVisible ? 'Show Teacher Answer Key & Rubric' : 'Hide Teacher Answer Key';
    }

    // Auto-save form fields to localStorage
    const INPUT_IDS = ['stud-name', 'stud-date', 'stud-period', 'q1-1', 'q1-2', 'q1-3', 'q2-1', 'q2-2', 'q2-3', 'q3-1', 'q3-2', 'q3-3', 'q4-1', 'q4-2', 'q4-3', 'q5-1', 'q5-2', 'q5-3'];

    INPUT_IDS.forEach(id => {
      const el = document.getElementById(id);
      if (el) {
        const saved = localStorage.getItem(`trouble_${id}`);
        if (saved) el.value = saved;
        el.addEventListener('input', () => {
          localStorage.setItem(`trouble_${id}`, el.value);
        });
      }
    });

    // Default Date to Today
    const dateInput = document.getElementById('stud-date');
    if (dateInput && !dateInput.value) {
      dateInput.value = new Date().toISOString().split('T')[0];
    }

    // Check for ?key=true in URL to auto-reveal teacher key
    const urlParams = new URLSearchParams(window.location.search);
    if (urlParams.get('key') === 'true') {
      toggleAnswerKey();
    }

    // Language handling
    function setLanguage(lang) {
      document.documentElement.lang = lang;
      document.documentElement.dir = lang === 'ar' ? 'rtl' : 'ltr';
      localStorage.setItem('activity_lang', lang);

      document.querySelectorAll('#lang-en, #lang-ar, #lang-uk').forEach(b => b.classList.remove('active'));
      const activeBtn = document.getElementById(`lang-${lang}`);
      if (activeBtn) activeBtn.classList.add('active');
    }

    const savedLang = localStorage.getItem('activity_lang') || 'en';
    setLanguage(savedLang);
  </script>
</body>
</html>
"""

with open(activity_html_path, "w", encoding="utf-8") as f:
    f.write(activity_html)

print(f"Generated activity.html -> {activity_html_path}")

# ==========================================
# 3. GENERATE NATIVE POWERPOINT DECK (PPTX)
# ==========================================
pptx_path = os.path.join(TARGET_DIR, "Divide_and_Conquer_Mastery.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

DARK_BG = RGBColor(4, 6, 8)
PANEL_BG = RGBColor(15, 23, 32)
GREEN = RGBColor(57, 255, 158)
CYAN = RGBColor(56, 189, 248)
AMBER = RGBColor(255, 180, 84)
RED = RGBColor(255, 92, 114)
WHITE = RGBColor(255, 255, 255)
TEXT_MUTED = RGBColor(148, 163, 184)
LINE_COLOR = RGBColor(30, 41, 59)

def add_slide_bg(slide):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BG
    bg_shape.line.fill.background()

def add_header(slide, badge_text, title_text, subtitle_text, badge_color=GREEN):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(4.5), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PANEL_BG
    badge.line.color.rgb = badge_color
    badge.line.width = Pt(1.5)
    tf_b = badge.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = badge_text
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = badge_color
    p_b.font.name = "Consolas"

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title_text
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = "Arial"

# Slide 1: Title
s1 = prs.slides.add_slide(blank_layout)
add_slide_bg(s1)
add_header(s1, "BEATTIE-NET // NETWORK ENGINEERING", "Divide & Conquer: CompTIA 6-Step Troubleshooting", "Mastering Systematic Fault Isolation & The Binary Search Diagnostic Method", GREEN)

card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.4), Inches(5.6), Inches(4.3))
card1.fill.solid(); card1.fill.fore_color.rgb = PANEL_BG; card1.line.color.rgb = GREEN; card1.line.width = Pt(2)
tf1 = card1.text_frame; tf1.word_wrap = True
p = tf1.paragraphs[0]; p.text = "01 // STOP GUESSING. START ISOLATING."; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = GREEN
p_b = tf1.add_paragraph(); p_b.text = "\nAmateur technicians guess and swap random parts. Enterprise network engineers and ethical hackers use formal deductive logic.\n\nToday you master the exact 6-step framework certified by CompTIA A+, Network+, and Security+."
p_b.font.size = Pt(13); p_b.font.color.rgb = WHITE

card2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.4), Inches(5.6), Inches(4.3))
card2.fill.solid(); card2.fill.fore_color.rgb = PANEL_BG; card2.line.color.rgb = CYAN; card2.line.width = Pt(2)
tf2 = card2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.text = "02 // THE BINARY SEARCH METHOD"; p2.font.size = Pt(16); p2.font.bold = True; p2.font.color.rgb = CYAN
p2_b = tf2.add_paragraph(); p2_b.text = "\nLinear troubleshooting through 7 OSI layers takes O(N) time.\n\nBy testing at the OSI Midpoint (Layer 3/4), you eliminate 50% of the entire problem domain in a single CLI command:\n\n• ping 127.0.0.1 (NIC)\n• ping Gateway (LAN)\n• ping 8.8.8.8 (WAN)\n• nslookup (DNS)"
p2_b.font.size = Pt(13); p2_b.font.color.rgb = WHITE

# Slide 2: 6 Steps Matrix
s2 = prs.slides.add_slide(blank_layout)
add_slide_bg(s2)
add_header(s2, "COMPTIA CORE CURRICULUM", "The 6-Step Troubleshooting Lifecycle", "Standard sequential methodology for CompTIA A+ / Network+ / Security+", CYAN)

steps = [
    ("1", "Identify the Problem", "Question user, review system logs, identify changes, duplicate symptom.", GREEN),
    ("2", "Establish Theory", "Question the obvious. Formulate probable cause hypothesis (Top/Bottom/Divide).", CYAN),
    ("3", "Test the Theory", "Execute targeted test. If confirmed -> proceed; if disproven -> form new theory.", AMBER),
    ("4", "Plan & Implement", "Develop step-by-step action plan, anticipate side effects, execute solution.", RED),
    ("5", "Verify & Prevent", "Verify full system functionality with end user; apply preventative hardening.", GREEN),
    ("6", "Document Everything", "Record findings, actions, and outcomes in enterprise ticketing & KB.", CYAN),
]

for idx, (num, title, desc, col) in enumerate(steps):
    r = idx // 3
    c = idx % 3
    x = Inches(0.8 + c * 3.9)
    y = Inches(2.4 + r * 2.3)
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.1))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL_BG; card.line.color.rgb = col; card.line.width = Pt(1.5)
    tf = card.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"STEP {num} // {title}"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = col
    p_d = tf.add_paragraph(); p_d.text = f"\n{desc}"; p_d.font.size = Pt(11); p_d.font.color.rgb = WHITE

# Slide 3: Midpoint Sequence
s3 = prs.slides.add_slide(blank_layout)
add_slide_bg(s3)
add_header(s3, "ENGINEERING WORKFLOW", "The 5-Step Midpoint Diagnostic Ladder", "Memorize these 5 sequential CLI commands to troubleshoot any network issue", AMBER)

ladder = [
    ("1", "ping 127.0.0.1", "Local Loopback Adapter", "Verifies local NIC driver and OS TCP/IP software stack.", CYAN),
    ("2", "ping 192.168.1.1", "Default Gateway (LAN Exit)", "Verifies patch cable, wall jack, switch port, and router interface (L1-L3).", GREEN),
    ("3", "ping 8.8.8.8", "External Public IP (WAN)", "Verifies ISP connection, NAT routing, and wide area internet path.", AMBER),
    ("4", "nslookup google.com", "DNS Name Resolution (L7)", "Verifies Port 53 DNS resolver, DNS caching server, and FQDN translation.", RED),
    ("5", "Test-NetConnection -Port 443", "Application Port & Service", "Verifies remote web service listening, TLS certificate, & stateful firewall.", CYAN),
]

for idx, (num, cmd, title, desc, col) in enumerate(ladder):
    y = Inches(2.3 + idx * 0.95)
    card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.85))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL_BG; card.line.color.rgb = col; card.line.width = Pt(1.5)
    tf = card.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"[{num}]  {cmd}  —  {title}: {desc}"
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE

# Slide 4: Case Studies
s4 = prs.slides.add_slide(blank_layout)
add_slide_bg(s4)
add_header(s4, "CTE REAL-WORLD LABS", "Hands-On Troubleshooting Case Studies", "Scenarios students will diagnose in today's companion lab activity", RED)

scens = [
    ("SCENARIO 1", "The Silent Wall Jack", "Physical Layer (L1/L2)", "Keystone jack miswire (pins 1-2 crossed). Wiremapping & punchdown repair.", GREEN),
    ("SCENARIO 2", "The 169.254 APIPA Storm", "Network Layer (L3 DHCP)", "DHCP scope exhaustion. Subnet expansion and lease time optimization.", AMBER),
    ("SCENARIO 3", "The DNS Dilemma", "Transport & DNS (L4/L7)", "8.8.8.8 pings OK, domain name fails. Stale static DNS resolver fix.", CYAN),
    ("SCENARIO 4", "Rogue SSH Exfiltration", "Incident Response", "Outbound Port 22 SSH tunnel alert. Endpoint quarantine and PID termination.", RED),
    ("SCENARIO 5", "The ARP Doppelgänger", "Enterprise Security", "Man-in-the-Middle ARP poisoning. Dynamic ARP Inspection & DHCP Snooping.", GREEN),
]

for idx, (tag, title, layer, desc, col) in enumerate(scens):
    y = Inches(2.3 + idx * 0.95)
    card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.85))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL_BG; card.line.color.rgb = col; card.line.width = Pt(1.5)
    tf = card.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{tag}: {title} [{layer}] — {desc}"
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE

prs.save(pptx_path)
print(f"Generated PowerPoint Presentation -> {pptx_path}")

# ==========================================
# 4. GENERATE MASTER INDEX.HTML LAUNCHPAD
# ==========================================
index_html_path = os.path.join(TARGET_DIR, "index.html")

index_html = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Beattie-Net // CompTIA 6-Step Troubleshooting Hub</title>
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;800&family=JetBrains+Mono:wght@400;700;800&family=Space+Grotesk:wght@600;700&display=swap" rel="stylesheet">
  <style>
    :root {
      --bg: #040608;
      --panel: #090e13;
      --panel-elevated: #101720;
      --line: #1c2b26;
      --green: #39ff9e;
      --cyan: #38bdf8;
      --amber: #ffb454;
      --text: #e2f4ea;
      --text-muted: #829a90;
    }
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body {
      background: var(--bg);
      color: var(--text);
      font-family: 'JetBrains Mono', monospace;
      min-height: 100vh;
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
      padding: 2rem;
      background-image: linear-gradient(rgba(57,255,158,0.04) 1px, transparent 1px), linear-gradient(90deg, rgba(57,255,158,0.04) 1px, transparent 1px);
      background-size: 36px 36px;
    }
    .hub-card {
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 12px;
      padding: 2.5rem 3rem;
      max-width: 900px;
      width: 100%;
      box-shadow: 0 20px 50px rgba(0,0,0,0.8);
      position: relative;
    }
    .badge {
      display: inline-block;
      font-size: 0.75rem;
      font-weight: 800;
      color: var(--green);
      background: rgba(57,255,158,0.08);
      border: 1px solid var(--green);
      padding: 0.3rem 0.8rem;
      border-radius: 4px;
      margin-bottom: 1rem;
    }
    h1 {
      font-family: 'Space Grotesk', sans-serif;
      font-size: 2.4rem;
      color: #fff;
      margin-bottom: 0.5rem;
    }
    p { color: var(--text-muted); font-size: 0.95rem; line-height: 1.6; margin-bottom: 2rem; }
    .grid { display: grid; grid-template-columns: 1fr 1fr; gap: 1.5rem; }
    .nav-box {
      background: var(--panel-elevated);
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 1.5rem;
      text-decoration: none;
      color: var(--text);
      display: flex;
      flex-direction: column;
      transition: all 0.2s ease;
    }
    .nav-box:hover {
      border-color: var(--green);
      transform: translateY(-2px);
      box-shadow: 0 10px 25px rgba(0,0,0,0.5);
    }
    .nav-box-title {
      font-family: 'Space Grotesk', sans-serif;
      font-size: 1.2rem;
      font-weight: 700;
      color: #fff;
      margin-bottom: 0.4rem;
    }
    .nav-box-desc { font-size: 0.82rem; color: var(--text-muted); line-height: 1.5; }
    .nav-box-tag { margin-top: auto; padding-top: 1rem; font-size: 0.75rem; color: var(--green); font-weight: 700; }
  </style>
</head>
<body>
  <div class="hub-card">
    <div class="badge">CTE CURRICULUM // MR. LINGSCH-BEATTIE</div>
    <h1>Divide &amp; Conquer Troubleshooting Hub</h1>
    <p>A.W. Beattie Career Center &middot; Network Engineering &amp; Cybersecurity. Master the CompTIA 6-Step Troubleshooting Framework &amp; Binary Search Logic.</p>

    <div class="grid">
      <a href="presentation.html" class="nav-box" style="border-left: 4px solid var(--green);">
        <div class="nav-box-title">📺 Slide Deck Presentation</div>
        <div class="nav-box-desc">Interactive 15-slide master lesson with live Midpoint Sandbox, presenter teaching notes, and keyboard controls.</div>
        <div class="nav-box-tag">Launch Slides (HTML) &rarr;</div>
      </a>

      <a href="activity.html" class="nav-box" style="border-left: 4px solid var(--cyan);">
        <div class="nav-box-title">⚡ Interactive Lab Activity</div>
        <div class="nav-box-desc">5 progressive troubleshooting tiers, interactive CLI terminal emulator, auto-saving forms, and printable worksheet PDF export.</div>
        <div class="nav-box-tag" style="color:var(--cyan);">Open Student Workbench &rarr;</div>
      </a>

      <a href="activity.html?key=true" class="nav-box" style="border-left: 4px solid var(--amber);">
        <div class="nav-box-title">🔑 Teacher Answer Key &amp; Rubric</div>
        <div class="nav-box-desc">Instructor review mode displaying complete answers, model flowchart logic, and grading criteria for all 5 scenarios.</div>
        <div class="nav-box-tag" style="color:var(--amber);">Open Teacher Rubric &rarr;</div>
      </a>

      <a href="Divide_and_Conquer_Mastery.pptx" download class="nav-box" style="border-left: 4px solid var(--text);">
        <div class="nav-box-title">📥 Native PowerPoint (.PPTX)</div>
        <div class="nav-box-desc">16:9 widescreen presentation deck ready for classroom projectors and offline PowerPoint presentations.</div>
        <div class="nav-box-tag" style="color:#fff;">Download PPTX File &rarr;</div>
      </a>
    </div>
  </div>
</body>
</html>
"""

with open(index_html_path, "w", encoding="utf-8") as f:
    f.write(index_html)

print(f"Generated index.html -> {index_html_path}")
print("ALL ASSETS GENERATED SUCCESSFULLY!")

